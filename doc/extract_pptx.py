from pptx import Presentation
import io

path = r"E:/antigravity/rag/doc/Ⅱ. 기술 및 기능_전달본_작성중_0430.pptx"
out_path = r"E:/antigravity/rag/doc/_pptx_dump.txt"

prs = Presentation(path)
lines = []
for i, slide in enumerate(prs.slides, 1):
    lines.append(f"\n========== SLIDE {i} ==========")
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = "".join(run.text for run in para.runs).strip()
                if text:
                    lines.append(text)
        if shape.shape_type == 19 or (hasattr(shape, "has_table") and shape.has_table):
            try:
                for row in shape.table.rows:
                    cells = [c.text.replace("\n", " \\n ").strip() for c in row.cells]
                    lines.append("[TABLE] " + " | ".join(cells))
            except Exception:
                pass
    # notes
    if slide.has_notes_slide:
        notes = slide.notes_slide.notes_text_frame.text.strip()
        if notes:
            lines.append("[NOTES] " + notes.replace("\n", " \\n "))

with io.open(out_path, "w", encoding="utf-8") as f:
    f.write("\n".join(lines))

print("done:", out_path, "slides:", len(prs.slides))
