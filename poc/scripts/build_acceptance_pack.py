"""고객 인수(acceptance) 샘플팩 생성툴 — 전 포맷 × 등급, 공개+합성 혼합.

목적(상용 운영 포장): 고객 폐쇄망 배포 후 parse→classify→gate 를 **실문서 포맷**으로 검증하는
자기완비 팩을 만든다. 분류 '정확도' 자체는 합성 천장(gold_real F1 0.26)이라 검증 목표가 아니다 —
팩이 실증하는 것은 (1) 파서가 전 포맷을 읽고 **숫자/표를 무손실 보존**하는지, (2) 고등급이
**under-분류(미탐)되지 않는지**(severity floor; 서빙은 안전방향 과분류라 정확일치가 아니라 floor 로 본다),
(3) 게이트 가시성(needs_review·경고)이다. 채점은 scripts/run_acceptance.py 가 한다.

시드(사용자 지시: 공개문서 혼합):
  - 합성 backbone : datasets/gold/classification_gold.jsonl  ([가상기업A] 안전 합성)
  - 공개 법근거   : datasets/gold_real/holdout_eval.clean.jsonl (legal_reference 동반)
  - 공개 S3 음성  : datasets/gold_real/... 없으면 합성 폴백
  실비밀 텍스트 0 — 모두 공개/합성. 소스 없으면 내장 폴백 시드로 그레이스풀 동작.

포맷: TXT · PDF · DOCX · XLSX · XLS · PPTX · HWPX (전부). HWP(바이너리)=파이썬 생성 불가 →
실 .hwp 픽스처(datasets/acceptance_seed/*.hwp)가 있으면 동봉, 없으면 스킵(HWPX 가 한국어 네이티브 대체).
reportlab(PDF)·xlwt(XLS)는 런타임/번들 불요 dev-only(scripts/requirements-docgen.txt; 팩이 미리 커밋되므로
번들은 이들이 불필요) → optional import, 없으면 해당 포맷 스킵. 렌더 팩은 커밋(airgap 결정적 배송);
전 포맷 재생성은 `make acceptance-pack-deps && make acceptance-pack`.

실행:  cd poc && python scripts/build_acceptance_pack.py
"""

from __future__ import annotations

import argparse
import json
import sys
import zipfile
from pathlib import Path
from xml.sax.saxutils import escape as _xml_escape

try:  # cp949 콘솔 안전 + pytest import 안전(재바인딩 대신 reconfigure)
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")  # type: ignore[union-attr]
except Exception:  # noqa: BLE001
    pass

_HERE = Path(__file__).resolve().parent
_ROOT = _HERE.parent  # poc/

GRADE_NAME = {"TS": "특급기밀", "S1": "1급비밀", "S2": "2급대외비", "S3": "3급공개"}
GRADES = ("TS", "S1", "S2", "S3")
HWPX_TEMPLATE = _ROOT / "tests" / "fixtures" / "sample.hwpx"

# 내장 폴백 시드 — 실 소스 데이터셋이 없어도(그레이스풀) 팩을 만들 수 있게. 공개(특허/판례 텍스처)와
# 합성([가상기업A])을 섞는다. 실비밀 아님 — 등급 신호는 명시 마킹 + 도메인 키워드로 준다.
_FALLBACK_SEEDS = [
    {"grade": "TS", "kind": "public", "legal_basis": "방위산업기술보호법 §3(방위산업기술)",
     "text": "본 기술문서는 국가핵심기술 및 방위산업기술에 해당하는 유도무기 표적탐지 신호처리 알고리즘의 "
             "설계 파라미터를 포함한다. 공개 시 국가안보에 중대한 영향을 미치므로 특급기밀로 관리한다."},
    {"grade": "TS", "kind": "synthetic", "legal_basis": "영업비밀(부정경쟁방지법 §2)",
     "text": "[가상기업A]의 최상위 영업비밀에 해당하는 반도체 EUV 노광 공정의 핵심 레시피와 장비 파라미터. "
             "경쟁사 대비 우위의 원천으로 유출 시 회복 불가능한 손해가 발생한다."},
    {"grade": "S1", "kind": "public", "legal_basis": "자본시장법 §174(미공개 중요정보)",
     "text": "당사 이사회의 미공개 중대 의사결정(대규모 인수합병 안건)과 관련한 내부 검토자료로, 공정공시 이전 "
             "단계의 중요정보이다. 미공개 상태에서의 유출·이용은 자본시장법 위반이다."},
    {"grade": "S1", "kind": "synthetic", "legal_basis": "영업비밀(핵심기술)",
     "text": "[가상기업A]의 고객 이탈예측 모델 설계서 및 내부 성능지표. 공개 모델 대비 정확도 우위를 담은 "
             "1급 비밀 자료로, 접근권한 없는 자의 열람을 금한다."},
    {"grade": "S2", "kind": "public", "legal_basis": "KOIPA 영업비밀 관리 가이드 §3.2.2",
     "text": "당사 금융 리스크 관리 모델의 내부 한도 설정 기준과 운영 절차를 담은 대외비 문서. 대외 공개는 "
             "제한되나 사내 관련 부서는 열람 가능한 2급 대외비로 분류한다."},
    {"grade": "S2", "kind": "synthetic", "legal_basis": "내부관리(대외비)",
     "text": "[가상기업A] 제조 공정 수율 개선 연구의 중간 보고서. 완성 전 단계의 내부 검토자료로 대외비로 관리한다."},
    {"grade": "S3", "kind": "public", "legal_basis": "공개 판결문(비공지성 상실)",
     "text": "이 문서는 이미 공간(公刊)된 대법원 판결문의 요지를 정리한 것으로, 누구나 열람 가능한 공개 정보이다. "
             "비밀관리성·비공지성이 없으므로 공개(S3) 등급이다."},
    {"grade": "S3", "kind": "public", "legal_basis": "사내 공지(비밀 아님)",
     "text": "전 임직원 대상 사내 복지제도 개편 안내. 민감정보나 영업비밀을 포함하지 않는 일반 공지문으로 공개 등급이다."},
]


def _norm(rec: dict) -> dict | None:
    """소스 레코드를 {grade, text, legal_basis} 로 정규화. 스키마 드리프트 흡수."""
    grade = rec.get("label") or rec.get("target") or rec.get("expected_grade") or rec.get("gold")
    text = rec.get("text") or rec.get("body") or rec.get("content")
    if grade not in GRADES or not text:
        return None
    return {
        "grade": grade,
        "text": " ".join(str(text).split())[:700],
        "legal_basis": rec.get("legal_reference") or rec.get("legal_basis") or "",
    }


def _load_jsonl(path: Path, kind: str, per_grade: int) -> list[dict]:
    """소스에서 등급별 per_grade 개를 결정적으로(doc_id 정렬) 추출."""
    if not path.exists():
        return []
    rows: list[dict] = []
    for line in path.read_text(encoding="utf-8").splitlines():
        line = line.strip()
        if line and not line.startswith("#"):
            try:
                rows.append(json.loads(line))
            except json.JSONDecodeError:
                continue
    rows.sort(key=lambda r: str(r.get("doc_id", "")))
    out: list[dict] = []
    seen: dict[str, int] = {g: 0 for g in GRADES}
    for r in rows:
        n = _norm(r)
        if not n or seen[n["grade"]] >= per_grade:
            continue
        n["kind"] = kind
        out.append(n)
        seen[n["grade"]] += 1
    return out


def load_seeds() -> list[dict]:
    """공개+합성 혼합 시드를 등급 균형으로 조립(결정적). 소스 없으면 내장 폴백."""
    seeds: list[dict] = []
    seeds += _load_jsonl(_ROOT / "datasets/gold_real/holdout_eval.clean.jsonl", "public", 1)
    seeds += _load_jsonl(_ROOT / "datasets/gold/classification_gold.jsonl", "synthetic", 1)
    have = {g: sum(1 for s in seeds if s["grade"] == g) for g in GRADES}
    # 폴백으로 등급별 최소 2건 보장(공개+합성 혼합 유지)
    for fb in _FALLBACK_SEEDS:
        if have.get(fb["grade"], 0) < 2:
            seeds.append(dict(fb))
            have[fb["grade"]] = have.get(fb["grade"], 0) + 1
    seeds.sort(key=lambda s: (GRADES.index(s["grade"]), s["kind"]))
    return seeds


def _numbers_for(idx: int) -> list[str]:
    """문서별 결정적·구별성 있는 숫자 3종(숫자 무손실 검증용 — 흔한 값 회피)."""
    base = 400000 + idx * 20731
    return [str(base), str(base + 137), str(base * 3 + 11)]


def _doc_spec(seed: dict, idx: int) -> dict:
    grade = seed["grade"]
    gname = GRADE_NAME[grade]
    doc_id = f"acc-{grade}-{idx:02d}"
    nums = _numbers_for(idx)
    title = f"[{gname}] {seed['kind']} 인수검증 표본 #{idx:02d}"
    header = f"문서번호: {doc_id}    보안등급: [{gname}]({grade})    법적근거: {seed.get('legal_basis') or 'N/A'}"
    footer = f"※ 본 문서는 [{gname}]({grade})으로 분류됨. 승인된 자 외 열람·유출 금지."
    # 통제 숫자표 — 표 포맷은 셀에, 산문 포맷은 인라인으로 동일 숫자를 심어 숫자무손실을 결정적으로 검증.
    table = [
        ["항목", "값"],
        ["문서번호", doc_id],
        ["개정판번호", nums[0]],
        ["대상연도금액", nums[1]],
        ["핵심수치", nums[2]],
    ]
    return {
        "doc_id": doc_id, "grade": grade, "gname": gname, "kind": seed["kind"],
        "legal_basis": seed.get("legal_basis") or "", "title": title, "header": header,
        "footer": footer, "body": seed["text"], "table": table, "numbers": nums,
    }


# ─────────────────────────────────────────────────────────────
# 포맷별 렌더러 (전부 한국어 + 숫자표). 반환 True=성공, False=라이브러리 부재로 스킵.
# ─────────────────────────────────────────────────────────────

def _body_lines(spec: dict) -> list[str]:
    return [spec["header"], spec["title"], spec["body"], spec["footer"]]


def render_txt(spec: dict, path: Path) -> bool:
    lines = _body_lines(spec)
    lines.append("")
    for row in spec["table"]:
        lines.append(" | ".join(row))
    path.write_text("\n".join(lines) + "\n", encoding="utf-8")
    return True


def render_docx(spec: dict, path: Path) -> bool:
    from docx import Document  # base dep
    doc = Document()
    doc.add_paragraph(spec["header"]).runs[0].bold = True
    doc.add_heading(spec["title"], level=1)
    for para in spec["body"].split("\n"):
        if para.strip():
            doc.add_paragraph(para.strip())
    t = doc.add_table(rows=len(spec["table"]), cols=2)
    t.style = "Table Grid"
    for r, row in enumerate(spec["table"]):
        t.rows[r].cells[0].text = row[0]
        t.rows[r].cells[1].text = row[1]
    doc.add_paragraph(spec["footer"]).runs[0].italic = True
    doc.save(str(path))
    return True


def render_pdf(spec: dict, path: Path) -> bool:
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
        from reportlab.lib.units import mm
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.cidfonts import UnicodeCIDFont
        from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table
    except ImportError:
        return False
    pdfmetrics.registerFont(UnicodeCIDFont("HYSMyeongJo-Medium"))  # 내장 CJK 폰트(외부 파일 불요)
    ss = getSampleStyleSheet()
    body = ParagraphStyle("kb", parent=ss["Normal"], fontName="HYSMyeongJo-Medium", fontSize=10, leading=16)
    head = ParagraphStyle("kh", parent=ss["Title"], fontName="HYSMyeongJo-Medium", fontSize=15)
    meta = ParagraphStyle("km", parent=ss["Normal"], fontName="HYSMyeongJo-Medium", fontSize=9)

    def esc(s: str) -> str:
        return s.replace("&", "&amp;").replace("<", "&lt;")

    story = [Paragraph(esc(spec["header"]), meta), Spacer(1, 3 * mm),
             Paragraph(esc(spec["title"]), head), Spacer(1, 3 * mm)]
    for para in spec["body"].split("\n"):
        if para.strip():
            story.append(Paragraph(esc(para.strip()), body))
            story.append(Spacer(1, 2 * mm))
    tbl = Table(spec["table"])
    story += [Spacer(1, 3 * mm), tbl, Spacer(1, 4 * mm), Paragraph(esc(spec["footer"]), meta)]
    SimpleDocTemplate(str(path), pagesize=A4, topMargin=18 * mm, bottomMargin=18 * mm).build(story)
    return True


def render_xlsx(spec: dict, path: Path) -> bool:
    from openpyxl import Workbook  # base dep
    wb = Workbook()
    ws = wb.active
    ws.title = "인수표본"
    for line in _body_lines(spec):
        ws.append([line])
    ws.append([])
    for row in spec["table"]:
        ws.append(row)
    wb.save(str(path))
    return True


def render_xls(spec: dict, path: Path) -> bool:
    try:
        import xlwt  # optional (pyproject 미선언)
    except ImportError:
        return False
    wb = xlwt.Workbook(encoding="utf-8")
    ws = wb.add_sheet("인수표본")
    r = 0
    for line in _body_lines(spec):
        ws.write(r, 0, line)
        r += 1
    r += 1
    for row in spec["table"]:
        ws.write(r, 0, row[0])
        ws.write(r, 1, row[1])
        r += 1
    wb.save(str(path))
    return True


def render_pptx(spec: dict, path: Path) -> bool:
    from pptx import Presentation  # base dep
    from pptx.util import Inches
    prs = Presentation()
    blank = prs.slide_layouts[6]
    s1 = prs.slides.add_slide(blank)
    tb = s1.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(5)).text_frame
    for line in _body_lines(spec):
        tb.add_paragraph().text = line
    s2 = prs.slides.add_slide(blank)
    rows, cols = len(spec["table"]), 2
    tbl = s2.shapes.add_table(rows, cols, Inches(0.5), Inches(0.5), Inches(8), Inches(0.4 * rows)).table
    for r, row in enumerate(spec["table"]):
        tbl.cell(r, 0).text = row[0]
        tbl.cell(r, 1).text = row[1]
    prs.save(str(path))
    return True


def render_hwpx(spec: dict, path: Path) -> bool:
    """HWPX 손수저작 — sample.hwpx 를 클론하고 section0.xml 만 본문+표로 교체.

    파서(extractor._hwpx_tables)는 any-prefix <tbl>/<tr>/<tc>/<t> 정규식이라 hhs: 접두어면 표 복원됨.
    rhwp.parse() 가 먼저 성공해야 하므로 템플릿 구조(hhs:section/p/run/t)에 밀착.
    """
    if not HWPX_TEMPLATE.exists():
        return False

    def t(s: str) -> str:
        return f"<hhs:t>{_xml_escape(s)}</hhs:t>"

    def para(s: str) -> str:
        return f"<hhs:p><hhs:run>{t(s)}</hhs:run></hhs:p>"

    rows_xml = "".join(
        "<hhs:tr>" + "".join(f"<hhs:tc>{t(c)}</hhs:tc>" for c in row) + "</hhs:tr>"
        for row in spec["table"]
    )
    section = (
        '<?xml version="1.0" encoding="UTF-8"?>'
        '<hhs:section xmlns:hhs="urn:schemas-hancom-com:hhs">'
        + "".join(para(ln) for ln in _body_lines(spec))
        + f"<hhs:tbl>{rows_xml}</hhs:tbl>"
        + "</hhs:section>"
    )
    src = zipfile.ZipFile(HWPX_TEMPLATE)
    path.parent.mkdir(parents=True, exist_ok=True)
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as out:
        for name in src.namelist():
            data = src.read(name)
            if name.endswith("section0.xml"):
                data = section.encode("utf-8")
            # mimetype 는 무압축이어야 안전 — 원본 스토어 방식 보존
            if name == "mimetype":
                out.writestr(name, data, compress_type=zipfile.ZIP_STORED)
            else:
                out.writestr(name, data)
    src.close()
    return True


_RENDERERS = {
    "txt": render_txt, "pdf": render_pdf, "docx": render_docx, "xlsx": render_xlsx,
    "xls": render_xls, "pptx": render_pptx, "hwpx": render_hwpx,
}
# 표 포맷(표 셀 추출·숫자무손실 강검증) 우선 배치 순서로 라운드로빈
_FORMAT_ORDER = ["docx", "xlsx", "hwpx", "pdf", "pptx", "xls", "txt"]


def build_pack(out_dir: Path) -> dict:
    docs_dir = out_dir / "docs"
    if docs_dir.exists():
        for f in docs_dir.iterdir():  # 생성 파일(acc-*)만 제거; 수기 real-* 실문서 픽스처는 보존
            if f.is_file() and f.name.startswith("acc-"):
                f.unlink()
    docs_dir.mkdir(parents=True, exist_ok=True)
    seeds = load_seeds()
    manifest_docs: list[dict] = []
    skipped: dict[str, int] = {}
    n = len(_FORMAT_ORDER)
    for idx, seed in enumerate(seeds, start=1):
        spec = _doc_spec(seed, idx)
        # 결정적 포맷: 문서 idx 로 고정 배정(라이브러리 부재 시에만 다음 포맷으로 폴백). 전 포맷 순회.
        start = (idx - 1) % n
        for fmt in _FORMAT_ORDER[start:] + _FORMAT_ORDER[:start]:
            renderer = _RENDERERS[fmt]
            fpath = docs_dir / f"{spec['doc_id']}.{fmt}"
            try:
                ok = renderer(spec, fpath)
            except Exception as exc:  # noqa: BLE001
                print(f"  [WARN] {spec['doc_id']} {fmt} 렌더 실패: {exc}")
                ok = False
            if ok:
                has_table = fmt in ("docx", "xlsx", "xls", "pptx", "hwpx")
                manifest_docs.append({
                    "doc_id": spec["doc_id"], "file": f"docs/{fpath.name}", "format": fmt,
                    "expected_grade": spec["grade"], "source_kind": spec["kind"],
                    "legal_basis": spec["legal_basis"],
                    "expected_numeric_tokens": spec["numbers"], "has_table": has_table,
                })
                break
            skipped[fmt] = skipped.get(fmt, 0) + 1
    # 수기 실문서 픽스처 병합 — 발주처 제공 실 .hwp/.hwpx(바이너리 파서 실검증)를 팩에 드롭인.
    # real_fixtures.json(out_dir)에 {file, expected_grade, expected_numeric_tokens, ...} 선언, docs/real-* 배치.
    real_fx = out_dir / "real_fixtures.json"
    n_real = 0
    if real_fx.exists():
        for d in json.loads(real_fx.read_text(encoding="utf-8")).get("docs", []):
            if (out_dir / d["file"]).exists():
                manifest_docs.append(d)
                n_real += 1
            else:
                print(f"  [WARN] real_fixtures 파일 없음(스킵): {d['file']}")

    manifest = {
        "pack": "lloydk-acceptance-pack",
        "note": ("고객 배포 인수검증용. 판정 규율: 정확 등급일치가 아니라 severity FLOOR(고등급 미탐=veto) "
                 "+ 숫자무손실 + 게이트 가시성. 소스=공개+합성 혼합, 실비밀 텍스트 0. HWP 바이너리는 생성불가."),
        "grades": GRADES,
        "docs": manifest_docs,
        "generated": len(manifest_docs) - n_real,
        "real_fixtures": n_real,
        "skipped_formats": skipped,
    }
    (out_dir / "expected_labels.json").write_text(
        json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")
    (out_dir / "README.md").write_text(
        "# Lloydk 인수(acceptance) 샘플팩\n\n"
        "고객 폐쇄망 배포 후 parse -> classify -> gate 를 실문서 포맷으로 검증하는 자기완비 팩.\n\n"
        "- 소스: 공개+합성 혼합. 실비밀 텍스트 0 (공개 판례/특허 텍스처 + [가상기업A] 합성).\n"
        "- 포맷: TXT/PDF/DOCX/XLSX/XLS/PPTX/HWPX (HWP 바이너리는 파이썬 생성 불가 -> HWPX 로 대체).\n"
        "- 판정: 정확 등급일치가 아니라 **severity floor**(고등급 미탐=veto) + 숫자 무손실 + 게이트 가시성.\n"
        "  서빙은 의도적 안전방향 과분류 -> 정확도로 합격/불합격하지 않는다.\n\n"
        "## 실행\n"
        "- 폐쇄망 번들(호스트): `API_KEY=<키> BASE_URL=http://localhost:8000 bash run_acceptance.sh`"
        " (run_acceptance.sh 는 번들 빌드시 자동 동봉)\n"
        "- 레포 보유(dev): `make acceptance-test` 또는 `python scripts/run_acceptance.py --mode inproc`\n\n"
        "재생성: `make acceptance-pack` (PDF/XLS 는 reportlab/xlwt 설치 시 포함).\n",
        encoding="utf-8")
    return manifest


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("--out", default="datasets/acceptance_pack", help="팩 출력 디렉토리(커밋 대상)")
    args = ap.parse_args()
    out_dir = _ROOT / args.out if not Path(args.out).is_absolute() else Path(args.out)
    m = build_pack(out_dir)
    from collections import Counter
    byfmt = Counter(d["format"] for d in m["docs"])
    bygrade = Counter(d["expected_grade"] for d in m["docs"])
    print(f"생성 완료: {len(m['docs'])} 문서 → {out_dir}")
    print(f"  포맷: {dict(byfmt)}")
    print(f"  등급: {dict(bygrade)}")
    if m["skipped_formats"]:
        print(f"  스킵(라이브러리 부재): {m['skipped_formats']} — reportlab(PDF)/xlwt(XLS) 는 optional")
    covered = set(byfmt)
    missing = [f for f in _FORMAT_ORDER if f not in covered]
    if missing:
        print(f"  [주의] 미커버 포맷: {missing} (해당 렌더 라이브러리 설치 후 재생성 권장)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
