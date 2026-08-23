"""m2_preprocess/extractor.py 분기 커버리지 보강.

외부 라이브러리 의존 분기는 라이브러리 미설치 시 graceful degrade 경로를 검증.
"""
from __future__ import annotations

from pathlib import Path

import pytest

from koipa.modules.m2_preprocess.extractor import ExtractResult, _pdf_tables_via_pdfplumber, extract


@pytest.fixture
def tmp_path_files(tmp_path: Path):
    """다양한 확장자의 임시 파일."""
    files = {}
    for ext, body in [
        ("txt", "테스트 텍스트 파일 본문\n두 번째 줄"),
        ("md", "# 마크다운 제목\n본문"),
        ("log", "2026-05-28 INFO 로그 라인"),
        ("csv", "col1,col2\nv1,v2"),
    ]:
        p = tmp_path / f"sample.{ext}"
        p.write_text(body, encoding="utf-8")
        files[ext] = p
    return files


class TestPlainExtraction:
    def test_txt_extraction(self, tmp_path_files):
        result = extract(tmp_path_files["txt"])
        assert isinstance(result, ExtractResult)
        assert result.method == "plain"
        assert result.quality == 1.0
        assert "테스트 텍스트 파일 본문" in result.text

    def test_md_extraction(self, tmp_path_files):
        result = extract(tmp_path_files["md"])
        assert result.method == "plain"
        assert "마크다운 제목" in result.text

    def test_log_extraction(self, tmp_path_files):
        result = extract(tmp_path_files["log"])
        assert result.method == "plain"
        assert "INFO" in result.text

    def test_csv_extraction(self, tmp_path_files):
        result = extract(tmp_path_files["csv"])
        assert result.method == "plain"
        assert "col1" in result.text

    def test_cp949_fallback(self, tmp_path: Path):
        """utf-8 실패 시 cp949 로 원문을 정확히 복원 (utf-16 모지바케 무음유실 회귀 방지).

        짝수 길이 순한글은 무BOM utf-16 이 예외 없이 임의 코드포인트로 '성공' 디코드하므로,
        폴백에서 utf-16 을 cp949 보다 먼저 시도하면 원문이 모지바케로 조용히 유실된다(silent FNR).
        폴백 순서가 cp949 우선(BOM 없을 때)이어야 원문이 그대로 복원된다.
        """
        original = "영업비밀"  # 4자 = cp949 8바이트(짝수) → utf-16 우선이면 모지바케로 깨짐
        p = tmp_path / "cp949.txt"
        p.write_bytes(original.encode("cp949"))
        result = extract(p)
        assert result.method == "plain"
        assert result.text == original, f"cp949 원문 복원 실패(모지바케 유실): {result.text!r}"
        assert any("cp949" in w for w in (result.warnings or []))


class TestUnsupported:
    def test_unsupported_extension(self, tmp_path: Path):
        p = tmp_path / "binary.xyz"
        p.write_bytes(b"\x00\x01\x02")
        result = extract(p)
        assert result.text == ""
        assert result.method == "plain"
        assert result.quality == 0.0
        assert "unsupported" in (result.error or "")

    def test_nonexistent_file(self, tmp_path: Path):
        """존재하지 않는 파일 — graceful 에러."""
        p = tmp_path / "does_not_exist.txt"
        # 존재하지 않으면 read 시점에서 에러 — extract가 안전하게 처리하는지
        try:
            result = extract(p)
            # 에러 또는 빈 결과 둘 다 OK
            assert isinstance(result, ExtractResult)
        except FileNotFoundError:
            # 명시적 에러도 OK (call site에서 처리)
            pass


class TestPdfFallback:
    """PDF 추출은 pdfminer.six → PyMuPDF fallback 경로 검증.

    실 PDF 없이도 _extract_pdf 호출 시 에러 메시지 정상 반환."""

    def test_pdf_nonexistent(self, tmp_path: Path):
        # PDF 확장자지만 실제 PDF 아닌 파일 → 라이브러리가 처리 시도 → 에러
        p = tmp_path / "fake.pdf"
        p.write_bytes(b"not a real pdf")
        result = extract(p)
        # 에러든 빈 결과든 ExtractResult 형태로 반환
        assert isinstance(result, ExtractResult)
        # quality가 0 (에러) 또는 매우 낮음
        assert result.quality <= 0.5

    def test_pdfplumber_tables_are_structured(self, tmp_path: Path, monkeypatch):
        import sys
        from types import SimpleNamespace

        class FakePage:
            def extract_tables(self):
                return [[["item", "amount"], ["secret", "12000"]]]

        class FakePdf:
            pages = [FakePage()]

            def __enter__(self):
                return self

            def __exit__(self, *args):
                return False

        monkeypatch.setitem(sys.modules, "pdfplumber", SimpleNamespace(open=lambda path: FakePdf()))

        p = tmp_path / "table.pdf"
        p.write_bytes(b"%PDF fake")
        tables, warnings = _pdf_tables_via_pdfplumber(p)
        assert warnings == ["pdf_tables_structured"]
        assert tables[0].rows == [["item", "amount"], ["secret", "12000"]]

    def test_pdf_table_extractor_unavailable_routes_review_signal(self, tmp_path: Path, monkeypatch):
        import sys

        from tests.test_document_ingestion import _minimal_pdf

        monkeypatch.setitem(sys.modules, "pdfplumber", None)

        p = tmp_path / "text_table_unknown.pdf"
        p.write_bytes(_minimal_pdf("Trade Secret ALD deposition table"))
        result = extract(p)

        assert result.method == "parser"
        assert result.text
        assert "pdf_table_extractor_unavailable" in result.warnings
        assert result.table_coverage == "incomplete"


class TestDocxFallback:
    def test_docx_nonexistent(self, tmp_path: Path):
        p = tmp_path / "fake.docx"
        p.write_bytes(b"not a real docx")
        result = extract(p)
        assert isinstance(result, ExtractResult)
        assert result.quality == 0.0
        assert result.error is not None

    def test_docx_table_is_structured(self, tmp_path: Path):
        docx = pytest.importorskip("docx")

        doc = docx.Document()
        table = doc.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "item"
        table.cell(0, 1).text = "amount"
        table.cell(1, 0).text = "secret cost"
        table.cell(1, 1).text = "12000"
        p = tmp_path / "table.docx"
        doc.save(str(p))

        result = extract(p)
        assert result.tables
        assert result.tables[0].rows == [["item", "amount"], ["secret cost", "12000"]]
        assert "secret cost | 12000" in result.text

    def test_docx_ooxml_content_control_table_fallback(self, tmp_path: Path, monkeypatch):
        import zipfile

        import docx

        class FakeDoc:
            paragraphs = []
            tables = []
            sections = []

        monkeypatch.setattr(docx, "Document", lambda path: FakeDoc())

        xml = (
            '<?xml version="1.0" encoding="UTF-8"?>'
            '<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
            "<w:body><w:tbl><w:tr>"
            "<w:tc><w:sdt><w:sdtContent><w:r><w:t>item</w:t></w:r></w:sdtContent></w:sdt></w:tc>"
            "<w:tc><w:sdt><w:sdtContent><w:r><w:t>amount</w:t></w:r></w:sdtContent></w:sdt></w:tc>"
            "</w:tr><w:tr>"
            "<w:tc><w:sdt><w:sdtContent><w:r><w:t>hidden weekly report</w:t></w:r></w:sdtContent></w:sdt></w:tc>"
            "<w:tc><w:r><w:t>260703</w:t></w:r></w:tc>"
            "</w:tr></w:tbl></w:body></w:document>"
        )
        p = tmp_path / "content_control_table.docx"
        with zipfile.ZipFile(p, "w", zipfile.ZIP_DEFLATED) as zf:
            zf.writestr("word/document.xml", xml)

        result = extract(p)
        assert "hidden weekly report | 260703" in result.text
        assert result.tables
        assert result.tables[0].rows == [["item", "amount"], ["hidden weekly report", "260703"]]
        assert "docx_ooxml_tables_extracted" in result.warnings


class TestHwpFallback:
    def test_hwp_nonexistent(self, tmp_path: Path):
        """HWP 라이브러리가 없거나 파일이 잘못된 경우 graceful degrade."""
        p = tmp_path / "fake.hwp"
        p.write_bytes(b"not a real hwp")
        result = extract(p)
        assert isinstance(result, ExtractResult)
        # rhwp 또는 parser 또는 에러 — 어떤 경우든 ExtractResult 형태
        assert result.method in ("rhwp", "parser")

    def test_hwpx_nonexistent(self, tmp_path: Path):
        p = tmp_path / "fake.hwpx"
        p.write_bytes(b"not a real hwpx")
        result = extract(p)
        assert isinstance(result, ExtractResult)


class TestExcelExtraction:
    """Excel(.xlsx/.xlsm/.xls) 추출 — 시트·표 셀을 텍스트로 (FUN-022 갭 보강 2026-06-05)."""

    def test_xlsx_extraction_real(self, tmp_path: Path):
        from openpyxl import Workbook

        wb = Workbook()
        ws = wb.active
        ws.title = "단가표"
        ws.append(["품목", "단가"])
        ws.append(["반도체 웨이퍼", 12000])
        ws2 = wb.create_sheet("비고")
        ws2.append(["기밀", "대외비"])
        p = tmp_path / "secret.xlsx"
        wb.save(str(p))

        result = extract(p)
        assert result.method == "openpyxl"
        assert result.quality > 0.5
        assert "단가표" in result.text          # 시트명 헤더
        assert "반도체 웨이퍼" in result.text     # 셀 값
        assert "12000" in result.text
        assert "대외비" in result.text           # 2번째 시트도
        assert "|" in result.text                # 표 셀 구분자

    def test_xlsx_table_comments_and_formula_are_preserved(self, tmp_path: Path):
        from openpyxl import Workbook
        from openpyxl.comments import Comment

        wb = Workbook()
        ws = wb.active
        ws.title = "Data"
        ws.append(["item", "amount"])
        ws.append(["secret cost", 12000])
        ws["A2"].comment = Comment("review this hidden note", "tester")
        ws["C2"] = "=B2*2"
        p = tmp_path / "structured.xlsx"
        wb.save(str(p))

        result = extract(p)
        assert result.tables
        assert result.tables[0].sheet == "Data"
        assert result.tables[0].rows[1][:2] == ["secret cost", "12000"]
        assert "[excel comment Data!A2] review this hidden note" in result.text
        assert "[excel formula Data!C2] =B2*2" in result.text
        assert "excel_comments_extracted" in result.warnings
        assert "excel_formulas_extracted" in result.warnings

    def test_xlsx_corrupt_graceful(self, tmp_path: Path):
        p = tmp_path / "fake.xlsx"
        p.write_bytes(b"not a real xlsx")
        result = extract(p)
        assert isinstance(result, ExtractResult)
        assert result.quality == 0.0
        assert result.error is not None

    def test_xls_graceful_degrade(self, tmp_path: Path):
        """구형 .xls — xlrd 있으면 파싱 시도(에러), 없으면 변환 안내. 둘 다 graceful."""
        p = tmp_path / "old.xls"
        p.write_bytes(b"\xd0\xcf\x11\xe0not a real xls")
        result = extract(p)
        assert isinstance(result, ExtractResult)
        assert result.quality == 0.0
        assert result.error is not None

    def test_xls_real_fixture(self):
        """구형 .xls 실파일(BIFF) 추출 — 다중시트 + 표. xlrd/픽스처 없으면 skip."""
        pytest.importorskip("xlrd", reason="구형 .xls 추출엔 xlrd 필요 (.[xls])")
        fx = Path(__file__).parent / "fixtures" / "sample.xls"
        if not fx.exists():
            pytest.skip("sample.xls 픽스처 없음")
        result = extract(fx)
        assert result.method == "xlrd", f"got {result.method}: {result.error}"
        assert result.quality > 0.5
        assert "영업비밀목록" in result.text       # 첫 시트명
        assert "ALD레시피" in result.text          # 셀 값
        assert "협력사" in result.text             # 둘째 시트도
        # ⚠ 이 문자열은 **브랜드 표기가 아니라 픽스처 내용**이다. sample.xls 는 바이너리라
        # 2026-08-12 koipa 리네임에서 치환 대상이 아니었고, 그 안의 셀 값은 "로이드케이"
        # 그대로다. 여기를 기관명으로 바꾸면 픽스처와 어긋나 테스트가 깨진다.
        assert "로이드케이" in result.text
        assert "|" in result.text                  # 표 셀 구분자

    def test_xlsx_merged_cell_vertical_propagation(self, tmp_path: Path):
        """세로 병합된 분류 라벨이 걸친 모든 행에 전파되는지 (표 의미 보존)."""
        from openpyxl import Workbook

        wb = Workbook()
        ws = wb.active
        ws.title = "등급표"
        ws.append(["분류", "문서", "등급"])         # row1
        ws["A2"] = "기밀"
        ws.merge_cells("A2:A4")                      # 세로 병합 A2:A4
        ws["B2"], ws["C2"] = "ALD레시피", "S1"
        ws["B3"], ws["C3"] = "원가구조", "S2"
        ws["B4"], ws["C4"] = "고객명단", "S2"
        p = tmp_path / "merged.xlsx"
        wb.save(str(p))

        text = extract(p).text
        # 병합 앵커값 '기밀'이 3개 데이터 행 각각에 함께 남는다.
        for doc in ("ALD레시피", "원가구조", "고객명단"):
            line = next(ln for ln in text.splitlines() if doc in ln)
            assert line.startswith("기밀 |"), f"세로병합 전파 실패: {line!r}"

    def test_xlsx_merged_cell_horizontal_no_repeat(self, tmp_path: Path):
        """가로 병합 헤더는 1회만 (중복 노이즈 회피)."""
        from openpyxl import Workbook

        wb = Workbook()
        ws = wb.active
        ws["A1"] = "2026년 분류표"
        ws.merge_cells("A1:C1")                      # 가로 병합
        ws.append([])
        ws.append(["분류", "문서", "등급"])
        p = tmp_path / "hmerge.xlsx"
        wb.save(str(p))

        text = extract(p).text
        assert text.count("2026년 분류표") == 1

    def test_xlsx_hidden_sheet_included(self, tmp_path: Path):
        """숨김 시트의 값도 추출하되 헤더에 상태 표기 (비밀 수치 누락 방지)."""
        from openpyxl import Workbook

        wb = Workbook()
        ws = wb.active
        ws.title = "공개"
        ws.append(["항목", "값"])
        hidden = wb.create_sheet("내부전용")
        hidden.sheet_state = "hidden"
        hidden.append(["진짜원가", 999999])
        p = tmp_path / "hidden.xlsx"
        wb.save(str(p))

        text = extract(p).text
        assert "내부전용 (hidden)" in text          # 상태 표기
        assert "999999" in text                     # 숨김 시트 값도 추출


class TestPptxExtraction:
    """PowerPoint(.pptx) 추출 — 슬라이드·도형 텍스트 (FUN-022 갭 보강 2026-06-05)."""

    def test_pptx_extraction_real(self, tmp_path: Path):
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "영업비밀 사업계획"
        tb = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
        tb.text_frame.text = "신제품 단가 12000원 대외비"
        p = tmp_path / "plan.pptx"
        prs.save(str(p))

        result = extract(p)
        assert result.method == "pptx"
        assert result.quality > 0.5
        assert "[slide 1]" in result.text
        assert "영업비밀 사업계획" in result.text     # 제목 도형
        assert "대외비" in result.text                # 텍스트박스

    def test_pptx_table_is_structured(self, tmp_path: Path):
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        shape = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(5), Inches(1))
        table = shape.table
        table.cell(0, 0).text = "item"
        table.cell(0, 1).text = "amount"
        table.cell(1, 0).text = "prototype"
        table.cell(1, 1).text = "9000"
        p = tmp_path / "table.pptx"
        prs.save(str(p))

        result = extract(p)
        assert result.tables
        assert result.tables[0].rows == [["item", "amount"], ["prototype", "9000"]]
        assert "prototype | 9000" in result.text

    def test_pptx_corrupt_graceful(self, tmp_path: Path):
        p = tmp_path / "fake.pptx"
        p.write_bytes(b"not a real pptx")
        result = extract(p)
        assert isinstance(result, ExtractResult)
        assert result.quality == 0.0
        assert result.error is not None


class TestDocExtraction:
    """구형 .doc(Word 97-2003) — antiword CLI 경로 (FUN-022 갭 보강 2026-06-05)."""

    def test_doc_without_antiword_graceful(self, tmp_path: Path, monkeypatch):
        import shutil

        monkeypatch.setattr(shutil, "which", lambda name: None)
        p = tmp_path / "old.doc"
        p.write_bytes(b"\xd0\xcf\x11\xe0fake doc")
        result = extract(p)
        assert result.method == "antiword"
        assert result.quality == 0.0
        assert "antiword" in (result.error or "") or "docx" in (result.error or "")

    def test_doc_routes_to_antiword(self, tmp_path: Path, monkeypatch):
        """.doc는 antiword 경로 — subprocess 모킹으로 추출 배선 검증."""
        import shutil
        import subprocess
        from types import SimpleNamespace

        monkeypatch.setattr(shutil, "which", lambda name: "antiword")
        monkeypatch.setattr(
            subprocess, "run",
            lambda args, **kw: SimpleNamespace(
                returncode=0, stdout="구형 워드 본문 영업비밀 단가표".encode("utf-8"), stderr=b""
            ),
        )
        p = tmp_path / "old.doc"
        p.write_bytes(b"\xd0\xcf\x11\xe0fake")
        result = extract(p)
        assert result.method == "antiword"
        assert "영업비밀" in result.text
        assert result.quality > 0.5


class TestPiiMaskingOnTables:
    """추출된 표 셀의 PII가 마스킹되는지 — 신규 포맷(xlsx/pptx) + 보안 결합 검증.

    감사(2026-06-05): 표/이미지 내 PII 마스킹 미검증. 표 셀도 추출 텍스트에 포함되므로
    rule 마스킹이 적용됨을 end-to-end로 고정(영업비밀 시스템 — PII 누출 방지).
    """

    def test_pii_in_xlsx_cell_is_masked(self, tmp_path: Path):
        from openpyxl import Workbook

        from koipa.modules.m2_preprocess.pii_masker import mask_pii

        wb = Workbook()
        ws = wb.active
        ws.append(["담당자", "주민번호", "연락처"])
        ws.append(["홍길동", "900101-1234567", "010-1234-5678"])
        p = tmp_path / "pii.xlsx"
        wb.save(str(p))

        text = extract(p).text
        assert "900101-1234567" in text          # 추출은 원문 그대로
        masked = mask_pii(text)
        assert "[RRN]" in masked.text             # 주민번호 마스킹
        assert "900101-1234567" not in masked.text
        assert "[PHONE]" in masked.text           # 전화번호도
        assert "010-1234-5678" not in masked.text

    def test_pii_in_pptx_is_masked(self, tmp_path: Path):
        from pptx import Presentation
        from pptx.util import Inches

        from koipa.modules.m2_preprocess.pii_masker import mask_pii

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        tb.text_frame.text = "계약 담당 이메일 secret@corp.co.kr 사업자 123-45-67890"
        p = tmp_path / "pii.pptx"
        prs.save(str(p))

        masked = mask_pii(extract(p).text)
        assert "[EMAIL]" in masked.text
        assert "secret@corp.co.kr" not in masked.text
        assert "[BIZNO]" in masked.text
        assert "123-45-67890" not in masked.text


class TestStringPath:
    def test_accepts_string_path(self, tmp_path: Path):
        """extract()는 str·Path 둘 다 받음."""
        p = tmp_path / "sample.txt"
        p.write_text("hello", encoding="utf-8")
        result = extract(str(p))
        assert result.text == "hello"
