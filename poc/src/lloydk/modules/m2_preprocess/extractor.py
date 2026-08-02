"""포맷별 텍스트 추출 라우터.

라이브러리 로드 실패는 graceful degrade — 라이브러리가 없는 환경에서도 .txt/.md는 처리.
extractor가 반환하는 ExtractResult는 추출 메서드·품질·OCR 사용 여부를 포함해
DB `documents.extraction_method/quality/ocr_used` 와 1:1.

OCR 엔진: Tesseract 5.x + 한국어팩 (kor.traineddata) — Apache 2.0.
  - 시스템 경로 자동 탐지 후 환경변수 TESSERACT_CMD 로 override 가능.
  - 폐쇄망 초기 설치: winget install UB-Mannheim.TesseractOCR 후 kor.traineddata 복사.
"""

from __future__ import annotations

import os
from dataclasses import dataclass, field
from pathlib import Path

# Tesseract 실행파일 경로 — 환경변수 > Windows 기본 경로 > PATH
_TESS_DEFAULT = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
TESSERACT_CMD = os.environ.get("TESSERACT_CMD") or (
    _TESS_DEFAULT if Path(_TESS_DEFAULT).exists() else "tesseract"
)

# poppler 경로 (pdf2image가 사용) — 환경변수 > 사용자 tools > None(PATH 탐색)
_POPPLER_CANDIDATES = [
    os.environ.get("POPPLER_PATH", ""),
    str(Path.home() / "tools" / "poppler" / "bin"),
    r"C:\Program Files\poppler\bin",
]
POPPLER_PATH: str | None = next(
    (p for p in _POPPLER_CANDIDATES if p and (Path(p) / "pdftoppm.exe").exists()),
    None,
)

# OCR DoS 가드 — 수백쪽 스캔 PDF 한 건이 pdf2image/Tesseract를 수십분~OOM으로 몰지
# 않도록 변환 페이지 수에 보수적 상한을 둔다. settings.ocr_max_pages가 있으면 그 값을,
# 없으면 모듈 상수(50)를 사용. 0 이하면 무제한으로 본다(명시적 opt-out).
_OCR_MAX_PAGES_DEFAULT = 50


SUPPORTED_FORMAT_GROUPS: dict[str, tuple[str, ...]] = {
    "plain": ("txt", "md", "log", "csv"),
    "hwp": ("hwp", "hwpx"),
    "word": ("docx", "doc"),
    "excel": ("xlsx", "xlsm", "xls"),
    "powerpoint": ("pptx", "pptm"),
    "pdf": ("pdf",),
    "image_ocr": ("jpg", "jpeg", "png", "tiff", "tif", "bmp", "webp"),
}
SUPPORTED_PARSE_EXTENSIONS: tuple[str, ...] = tuple(
    ext for group in SUPPORTED_FORMAT_GROUPS.values() for ext in group
)


def _ocr_max_pages() -> int:
    try:
        from lloydk.config import settings  # noqa: PLC0415

        v = getattr(settings, "ocr_max_pages", None)
        if v is not None:
            return int(v)
    except Exception:  # noqa: BLE001
        pass
    return _OCR_MAX_PAGES_DEFAULT


# ── zip bomb(압축폭탄) 가드 ────────────────────────────────────────────────────
# DOCX·HWPX·xlsx·pptx 는 모두 zip 컨테이너다. max_upload_mb 는 '압축된' 업로드 크기만 제한하므로,
# 수십 KB 파일이 GB급으로 팽창하는 압축폭탄을 막지 못한다(추출 단계 OOM/DoS). infolist()의 비압축
# 크기(file_size)·압축비로 **실제 decompress 전에** 사전 차단한다(central directory 메타만 읽음).
_ZIP_MAX_TOTAL_UNCOMPRESSED = 512 * 1024 * 1024   # 총 비압축 상한(512MB) — 정상 오피스문서 여유
_ZIP_MAX_COMPRESSION_RATIO = 200                   # 멤버별 압축비 상한(비정상 팽창=폭탄 시그니처)


def _guard_zip_bomb(zf, *, source: str = "archive") -> None:
    """zip 멤버 메타데이터로 압축폭탄을 사전 차단 — 실제 read 전에 검사(decompress 0).

    zipfile.read()는 선언된 file_size 를 초과하면 CRC/크기 검증에서 실패하므로, central directory
    의 file_size 합은 실제 팽창량의 신뢰 가능한 상한이다. 총량·멤버 압축비 둘 다 본다.
    """
    total = 0
    for info in zf.infolist():
        size = int(getattr(info, "file_size", 0) or 0)
        comp = int(getattr(info, "compress_size", 0) or 0) or 1
        total += size
        if total > _ZIP_MAX_TOTAL_UNCOMPRESSED:
            raise ValueError(
                f"zip decompression bomb suspected ({source}): uncompressed total "
                f"> {_ZIP_MAX_TOTAL_UNCOMPRESSED // (1024 * 1024)}MB"
            )
        if size and (size / comp) > _ZIP_MAX_COMPRESSION_RATIO:
            raise ValueError(
                f"zip member compression ratio too high ({source}: {info.filename}): "
                f"{size}/{comp} > {_ZIP_MAX_COMPRESSION_RATIO}"
            )


@dataclass
class ExtractedTable:
    source: str
    name: str
    rows: list[list[str]]
    sheet: str | None = None
    page: int | None = None
    table_index: int | None = None


@dataclass
class ExtractResult:
    text: str
    method: str           # parser/rhwp/ocr/ocr_llm/libreoffice/plain
    quality: float        # 0.0~1.0 (추정)
    ocr_used: bool = False
    pages: int | None = None
    error: str | None = None
    # 표 커버리지 신호 — 원문에 표가 있는데 셀 텍스트가 추출본에 빠졌을 때 "incomplete".
    # HWP/HWPX에서 rhwp가 표·글상자 셀을 흘리면 본문만 분류기에 들어가 표 속 영업비밀이
    # '조용히' 미탐(FNR)되는 것을 가시화한다. None=해당없음/회수됨/판정불가. 메서드 고정
    # quality(0.95 등)는 '무엇을 놓쳤나'가 아니라 '어떤 파서를 썼나'만 반영해 부분손실을
    # 가리므로, 이 신호로 검수 라우팅을 발동한다(등급 불변 — FNR-safe).
    table_coverage: str | None = None
    tables: list[ExtractedTable] = field(default_factory=list)
    warnings: list[str] = field(default_factory=list)


def extract(path: str | Path) -> ExtractResult:
    p = Path(path)
    suffix = p.suffix.lower().lstrip(".")
    if suffix in SUPPORTED_FORMAT_GROUPS["plain"]:
        return _extract_plain(p)
    if suffix in SUPPORTED_FORMAT_GROUPS["hwp"]:
        return _extract_hwp(p)
    if suffix == "docx":
        return _extract_docx(p)
    if suffix == "doc":
        return _extract_doc(p)
    if suffix in SUPPORTED_FORMAT_GROUPS["excel"]:
        return _extract_excel(p)
    if suffix in SUPPORTED_FORMAT_GROUPS["powerpoint"]:
        return _extract_pptx(p)
    if suffix == "pdf":
        return _extract_pdf(p)
    if suffix in SUPPORTED_FORMAT_GROUPS["image_ocr"]:
        return _extract_image_ocr(p)
    return ExtractResult(text="", method="plain", quality=0.0, error=f"unsupported: {suffix}")


def _extract_plain(p: Path) -> ExtractResult:
    try:
        text = p.read_text(encoding="utf-8")
        return ExtractResult(text=text, method="plain", quality=1.0)
    except UnicodeDecodeError:
        pass
    raw = p.read_bytes()
    warnings: list[str] = []
    # 인코딩 폴백 순서 주의(무음 유실 방지): bare 'utf-16'(무BOM)은 짝수 길이 CP949/EUC-KR
    # 한국어 바이트열을 예외 없이 임의 코드포인트로 '성공' 디코드한다. 과거엔 utf-16 을 cp949 보다
    # 먼저 시도해 레거시 한국어 평문(.txt/.csv/.log/.md·CSV 내보내기·로그)이 모지바케로 무음 유실됐다
    # (품질 0.9 로 정상 추출 위장 → 한국어 신호 상실 → silent FNR). BOM 이 있을 때만 utf-16 을 쓰고,
    # 그 외에는 cp949/euc-kr 를 먼저 시도하고 utf-16-le/be 는 최후 폴백으로 둔다(euc-kr⊂cp949).
    if raw[:2] in (b"\xff\xfe", b"\xfe\xff"):
        candidates = ("utf-16", "cp949", "euc-kr")
    else:
        candidates = ("cp949", "euc-kr", "utf-16-le", "utf-16-be")
    for enc in candidates:
        try:
            text = raw.decode(enc)
            warnings.append(f"plain_decoded_as_{enc}")
            return ExtractResult(text=text, method="plain", quality=0.9, warnings=warnings)
        except UnicodeDecodeError:
            continue
    text = raw.decode("utf-8", errors="replace")
    replacement_count = text.count("\ufffd")
    warnings.append(f"plain_decode_replaced_chars:{replacement_count}")
    quality = max(0.3, 1.0 - min(0.7, replacement_count / max(1, len(text))))
    return ExtractResult(text=text, method="plain", quality=round(quality, 4), warnings=warnings)


def _cell_text(value) -> str:
    if value is None:
        return ""
    return str(value).strip()


def _xls_cell_text(value, ctype, datemode) -> str:
    """xlrd 셀을 타입 인지 문자열로 — 정수는 '.0' 없이, 날짜는 serial 대신 datetime 으로.

    xlrd 는 모든 수치/날짜를 float 로 돌려줘 str() 만 하면 1234→'1234.0', 엑셀 날짜(serial 45306)→
    '45306.0'(의미 소실)로 열화된다(openpyxl 경로는 int/datetime 보존). 숫자무손실 위해 교정.
    """
    import xlrd  # noqa: PLC0415
    if ctype == xlrd.XL_CELL_NUMBER:
        try:
            if float(value).is_integer():
                return str(int(value))
        except (TypeError, ValueError, OverflowError):
            pass
        return str(value)
    if ctype == xlrd.XL_CELL_DATE:
        try:
            dtv = xlrd.xldate.xldate_as_datetime(value, datemode)
            return dtv.date().isoformat() if (dtv.hour == dtv.minute == dtv.second == 0) else dtv.isoformat(sep=" ")
        except Exception:  # noqa: BLE001
            return str(value)
    return _cell_text(value)


def _trim_trailing_empty(row: list[str]) -> list[str]:
    while row and not row[-1]:
        row.pop()
    return row


def _row_has_text(row: list[str]) -> bool:
    return any(c.strip() for c in row)


def _table_text_line(row: list[str]) -> str:
    return " | ".join(c for c in row if c.strip())


def _tables_to_text(tables: list[ExtractedTable]) -> str:
    lines: list[str] = []
    for table in tables:
        for row in table.rows:
            line = _table_text_line(row)
            if line:
                lines.append(line)
    return "\n".join(lines)


def _append_if_missing(text: str, addition: str) -> str:
    if not addition.strip():
        return text or ""
    norm = "".join((text or "").split())
    missing = []
    for line in addition.splitlines():
        if line.strip() and "".join(line.split()) not in norm:
            missing.append(line)
    if not missing:
        return text or ""
    return ((text or "").rstrip() + "\n" + "\n".join(missing)).strip()


def _warn_once(warnings: list[str], warning: str) -> None:
    if warning not in warnings:
        warnings.append(warning)


def _zip_names(path: Path) -> list[str]:
    import zipfile

    try:
        with zipfile.ZipFile(path) as zf:
            return zf.namelist()
    except Exception:  # noqa: BLE001
        return []


_PACKAGE_IMAGE_SUFFIXES = (".png", ".jpg", ".jpeg", ".tif", ".tiff", ".bmp", ".webp")
_PACKAGE_OCR_LIMIT = 30


def _ocr_package_images(
    p: Path,
    *,
    prefix: str,
    label: str,
    warnings: list[str],
) -> list[str]:
    import io
    import zipfile

    names = [
        n for n in _zip_names(p)
        if n.startswith(prefix) and n.lower().endswith(_PACKAGE_IMAGE_SUFFIXES)
    ]
    if not names:
        return []
    try:
        from PIL import Image  # type: ignore
        import pytesseract  # type: ignore  # noqa: F401
    except ImportError:
        _warn_once(warnings, f"{label}_media_ocr_unavailable")
        return []

    parts: list[str] = []
    try:
        with zipfile.ZipFile(p) as zf:
            _guard_zip_bomb(zf, source="office-zip")
            for name in names[:_PACKAGE_OCR_LIMIT]:
                try:
                    with Image.open(io.BytesIO(zf.read(name))) as img:
                        text = _tess_image(img)
                except Exception:  # noqa: BLE001
                    continue
                if text and text.strip():
                    parts.append(f"[{label} image OCR {name}]\n{text.strip()}")
                    try:
                        warnings.remove(f"{label}_media_not_ocrd")
                    except ValueError:
                        pass
                    _warn_once(warnings, f"{label}_media_ocr_extracted")
    except Exception as exc:  # noqa: BLE001
        _warn_once(warnings, f"{label}_media_ocr_error:{type(exc).__name__}")
        return parts

    if len(names) > _PACKAGE_OCR_LIMIT:
        _warn_once(warnings, f"{label}_media_ocr_truncated")
    if not parts:
        _warn_once(warnings, f"{label}_media_ocr_empty")
    return parts


def _xml_text(xml: str) -> str:
    import html as _html
    import re as _re

    chunks = _re.findall(r"<(?:\w+:)?t\b[^>]*>(.*?)</(?:\w+:)?t>", xml, _re.S)
    text = "\n".join(_html.unescape(_re.sub(r"<[^>]+>", "", c)).strip() for c in chunks)
    return "\n".join(line for line in text.splitlines() if line.strip())


def _html_to_text(raw: str) -> str:
    import html as _html
    import re as _re

    raw = _re.sub(r"<(script|style)[^>]*>.*?</\1>", " ", raw, flags=_re.S | _re.I)
    text = _re.sub(r"<[^>]+>", " ", raw)
    text = _html.unescape(text)
    text = _re.sub(r"[ \t\u00a0]+", " ", text)
    text = _re.sub(r"\s*\n\s*", "\n", text).strip()
    return text


def _html_tables(raw: str, *, source: str, name_prefix: str) -> list[ExtractedTable]:
    import html as _html
    import re as _re

    tables: list[ExtractedTable] = []
    for tidx, table in enumerate(
        _re.findall(r"<table\b[^>]*>(.*?)</table>", raw, _re.S | _re.I),
        start=1,
    ):
        rows: list[list[str]] = []
        for tr in _re.findall(r"<tr\b[^>]*>(.*?)</tr>", table, _re.S | _re.I):
            cells: list[str] = []
            for cell in _re.findall(r"<t[dh]\b[^>]*>(.*?)</t[dh]>", tr, _re.S | _re.I):
                text = _html_to_text(_html.unescape(cell))
                cells.append(text)
            cells = _trim_trailing_empty(cells)
            if _row_has_text(cells):
                rows.append(cells)
        if rows:
            tables.append(
                ExtractedTable(
                    source=source,
                    name=f"{name_prefix} table {tidx}",
                    rows=rows,
                    table_index=tidx,
                )
            )
    return tables


def _extract_hwp(p: Path) -> ExtractResult:
    """HWP5/HWPX 추출 — rhwp-python(Rust PyO3, HWP5+HWPX 통합 파서, extra ``[hwp]``).

    rhwp가 .hwp(바이너리)·.hwpx 둘 다 파싱한다. 미설치/파싱 실패는 graceful degrade.
    구 pyhwp/hwp5 폴백 제거(2026-06-05): 설치돼 있던 hwp5 0.1.0은 무관 패키지였고
    ``hwp5.extract_text`` API도 부재해 항상 실패하던 죽은 경로였다.

    알려진 한계(실측 2026-06-27, 공개 보도자료·서식 .hwp 대상, rhwp 0.5.1·0.8.0 동일):
      본문 단락은 정상 추출하나 **표·글상자 셀 텍스트를 거의 전부 흘린다**
      (noori.hwp 표셀 29/31 누락, basicsReport.hwp 23/23 누락). 영업비밀 문서는
      등급·원가·담당 정보가 표에 들어가는 일이 많아 '조용한 표 미추출'이 곧 미탐이다.

    표 보강(2026-08-01, unhwp/MIT — 구 pyhwp/AGPL 대체): **.hwp 에서 표 누락이
      의심될 때만**(table_coverage=="incomplete") unhwp 로 표 셀을 회수해 rhwp
      본문에 **합집합으로** 덧붙인다(_hwp_tables_via_unhwp). 세 가지가 모두 실측
      근거다(8문서×3경로 A/B):
        - .hwp 는 rhwp 가 표를 통째로 흘린다 → unhwp 가 593행 회수, 숫자 7/7 보존 + 21개 추가.
        - .hwpx 는 반대다. rhwp 가 이미 잘 뽑고 unhwp 가 표를 통째로 누락한 사례가 있다
          (acc-S1-03.hwpx: 표 5행→0행, 인수 기대숫자 3/3→0/3) → **.hwpx 는 건드리지 않는다**.
        - 따라서 교체가 아니라 합집합이다. unhwp 는 rhwp 의 상위집합이 아니다.
      보강에 성공해도 table_coverage 는 "incomplete" 로 유지한다 — 표 내용을 분류기에
      넣는 이득(미탐 감소)과 검수 라우팅 완화는 별개 결정이며, 후자는 현장 데이터
      확보 후 판단한다(확실한 사람 검수를 신규 파서 신뢰와 맞바꾸지 않는다).
      미설치·실패는 graceful — rhwp 결과를 그대로 쓴다.
    """
    try:
        import rhwp  # type: ignore
    except ImportError as exc:
        return ExtractResult(
            text="", method="rhwp", quality=0.0,
            error=f"HWP 추출에는 rhwp-python 필요 (pip install '.[hwp]'): {exc}",
        )
    try:
        doc = rhwp.parse(str(p))
        text = doc.extract_text()
    except Exception as exc:  # noqa: BLE001
        return ExtractResult(text="", method="rhwp", quality=0.0, error=str(exc))

    warnings: list[str] = []
    tables: list[ExtractedTable] = []
    is_hwpx = p.suffix.lower() == ".hwpx"
    if is_hwpx:
        # .hwpx 는 원본 zip 의 section XML 에서 셀을 직접 뽑는다(정밀). unhwp 보강 대상 아님.
        tables = _hwpx_tables(p.read_bytes(), source="hwpx")
        if tables:
            text = _append_if_missing(text or "", _tables_to_text(tables))
            _warn_once(warnings, "hwpx_tables_structured")

    if not text or not text.strip():
        # 빈 추출 = 성공 아님. 표/글상자 전용 HWP일 가능성을 경고로 surface.
        return ExtractResult(
            text="", method="rhwp", quality=0.0,
            error="rhwp가 본문 텍스트를 추출하지 못함 — 표·글상자 전용 문서일 수 있음. "
                  "표 셀 회수 보강에는 unhwp 필요: pip install '.[hwp-tables]'.",
        )
    # 본문은 뽑혔지만 표 셀이 빠졌을 수 있다(rhwp의 조용한 표 미추출) — 커버리지 판정.
    coverage = _hwp_table_coverage(p, doc, text, warnings)

    # [FNR] 판정 불가(None + hwp_table_check_unavailable)도 보강 대상이다.
    # 종전엔 보강이 coverage=="incomplete" 안에만 있어서, 표 검출기(to_hwpx_bytes)가
    # 실패하면 "표가 있는지 모른다"는 이유로 보강도 검수 라우팅도 건너뛰었다 — fail-open.
    #
    # 실측(공고문 .hwp · 2026-08-02): rhwp-python 0.5.1 처럼 to_hwpx_bytes 가 없는
    # 환경에서 본문 10,439자·표 0개로 자동확정됐다. 같은 파일이 0.8.1(배포 이미지)에서는
    # 검출기가 동작해 46,473자·표 47개·셀 1,815개로 정상 회수·검수 라우팅된다.
    # 즉 이 분기는 배포본의 현재 동작이 아니라 rhwp 버전·파일에 따라 열리는 잠재 경로다.
    # 그래도 닫는 이유: 검출기가 "모른다"고 답할 때 무음 통과시키면 표 속 비밀이 그대로
    # 미탐이 된다. 모름을 안전으로 읽지 않는다.
    unknown = coverage is None and "hwp_table_check_unavailable" in (warnings or [])
    if coverage == "incomplete" or unknown:
        if coverage == "incomplete":
            _warn_once(warnings, "hwp_table_cells_may_be_missing")
        # .hwp 한정 표 보강. 합집합이라 rhwp 가 이미 잡은 내용은 중복되지 않는다.
        # coverage 는 그대로 "incomplete" 로 둔다 → 검수 라우팅 유지(설계 의도).
        if not is_hwpx:
            recovered, unhwp_tables = _hwp_tables_via_unhwp(p)
            if recovered:
                merged = _append_if_missing(text, recovered)
                if len(merged) > len(text):
                    text = merged
                    tables = tables + unhwp_tables
                    _warn_once(warnings, "hwp_tables_recovered_by_unhwp")
                    # 회수됐다는 것은 rhwp 가 실제로 표를 흘렸다는 뜻 — 판정 불가였더라도
                    # 이제는 유실이 확인된 것이므로 incomplete 로 확정해 검수로 보낸다.
                    coverage = "incomplete"
                    _warn_once(warnings, "hwp_table_cells_may_be_missing")
            elif unknown:
                # 표 유무도 모르고 회수도 못 했다 = 유실 없음을 입증할 수 없다.
                # 무음 통과 대신 검수로 보낸다(등급은 불변 — FNR-safe).
                _warn_once(warnings, "hwp_table_cells_may_be_missing")
                coverage = "incomplete"
    return ExtractResult(
        text=text,
        method="rhwp",
        quality=0.95,
        table_coverage=coverage,
        tables=tables,
        warnings=warnings,
    )


def _hwp_table_coverage(p: Path, doc, text: str, warnings: list[str] | None = None) -> str | None:
    """rhwp가 본문은 뽑았으나 표 셀을 흘렸는지 판정 → "incomplete" | None.

    표 속 등급·원가·담당·임원보상 같은 영업비밀이 본문에 없으면 분류기가 못 보고
    저등급으로 '조용히' 미탐한다. 메서드 고정 quality(0.95)로는 안 잡히는 사각.

    - **.hwpx**: 원본 zip의 section XML에서 표 셀 텍스트(<...:t>)를 뽑아 추출본과 대조.
      셀 텍스트가 추출본에 없으면 "incomplete"(정밀 — rhwp가 회수했으면 flag 안 함).
    - **.hwp(바이너리)**: 셀 텍스트는 rhwp HWPX 변환에도 빠지므로 대조 불가. 표 *구조*
      존재만 보고 "incomplete". 보수적(FNR-safe).
      이 판정이 곧 unhwp 표 보강의 트리거다 — 보강에 성공해도 "incomplete" 는
      유지한다(검수 라우팅 완화는 별개 결정). 즉 이 값은 "표 셀 미회수"가 아니라
      "rhwp 단독으로는 표를 못 봤다"는 뜻이다.

    판정 실패·예외·표 없음은 None(graceful) — 과라우팅을 피한다.
    """
    try:
        suffix = p.suffix.lower()
        if suffix == ".hwpx":
            if _hwpx_uncaptured_table_cells(p.read_bytes(), text or ""):
                return "incomplete"
            return None
        if suffix == ".hwp":
            try:
                hwpx_bytes = doc.to_hwpx_bytes()
            except Exception:  # noqa: BLE001
                # 유일한 표-존재 검출기(to_hwpx_bytes)가 이 .hwp에서 실패 — 표 유무 판정 불가.
                # 조용히 None으로 삼키면 표 속 비밀 미탐이 무음이 된다. 가시화(등급 불변).
                if warnings is not None:
                    _warn_once(warnings, "hwp_table_check_unavailable")
                return None
            if _hwpx_bytes_has_table(hwpx_bytes):
                return "incomplete"
            return None
    except Exception:  # noqa: BLE001
        return None
    return None


def _hwpx_uncaptured_table_cells(hwpx_bytes: bytes, extracted_text: str) -> bool:
    """HWPX(zip)의 표 셀 텍스트 중 추출본에 없는 게 하나라도 있으면 True.

    section XML의 <...:tbl> 블록 안 <...:t> 텍스트를 모아, 공백 제거 후 추출본(공백 제거)에
    부분문자열로 존재하는지 본다. 접두 네임스페이스 무관(hp:/hhs:/무접두 모두 매칭).
    셀 텍스트가 이미 본문에 있으면(회수됨/반복) flag 안 함 — 오탐 최소화.
    """
    import html as _html
    import io
    import re as _re
    import zipfile

    norm = _re.sub(r"\s+", "", extracted_text or "")
    try:
        z = zipfile.ZipFile(io.BytesIO(hwpx_bytes))
        _guard_zip_bomb(z, source="hwpx")
    except Exception:  # noqa: BLE001
        return False
    names = [n for n in z.namelist() if n.lower().endswith(".xml") and "section" in n.lower()]
    if not names:
        names = [n for n in z.namelist() if n.lower().endswith(".xml")]
    for n in names:
        try:
            xml = z.read(n).decode("utf-8", errors="replace")
        except Exception:  # noqa: BLE001
            continue
        for tbl in _re.findall(r"<(?:\w+:)?tbl\b[^>]*>(.*?)</(?:\w+:)?tbl>", xml, _re.S):
            for cell in _re.findall(r"<(?:\w+:)?t\b[^>]*>(.*?)</(?:\w+:)?t>", tbl, _re.S):
                c = _html.unescape(_re.sub(r"<[^>]+>", "", cell)).strip()
                if len(c) >= 2 and _re.sub(r"\s+", "", c) not in norm:
                    return True
    return False


def _hwpx_tables(hwpx_bytes: bytes, *, source: str = "hwpx") -> list[ExtractedTable]:
    import html as _html
    import io
    import re as _re
    import zipfile

    try:
        z = zipfile.ZipFile(io.BytesIO(hwpx_bytes))
        _guard_zip_bomb(z, source="hwpx")
    except Exception:  # noqa: BLE001
        return []
    names = [n for n in z.namelist() if n.lower().endswith(".xml") and "section" in n.lower()]
    if not names:
        names = [n for n in z.namelist() if n.lower().endswith(".xml")]

    out: list[ExtractedTable] = []
    for name in names:
        try:
            xml = z.read(name).decode("utf-8", errors="replace")
        except Exception:  # noqa: BLE001
            continue
        for tbl in _re.findall(r"<(?:\w+:)?tbl\b[^>]*>(.*?)</(?:\w+:)?tbl>", xml, _re.S):
            rows: list[list[str]] = []
            tr_blocks = _re.findall(r"<(?:\w+:)?tr\b[^>]*>(.*?)</(?:\w+:)?tr>", tbl, _re.S)
            if not tr_blocks:
                tr_blocks = [tbl]
            for tr in tr_blocks:
                row: list[str] = []
                for tc in _re.findall(r"<(?:\w+:)?tc\b[^>]*>(.*?)</(?:\w+:)?tc>", tr, _re.S):
                    texts = []
                    for t in _re.findall(r"<(?:\w+:)?t\b[^>]*>(.*?)</(?:\w+:)?t>", tc, _re.S):
                        cell = _html.unescape(_re.sub(r"<[^>]+>", "", t)).strip()
                        if cell:
                            texts.append(cell)
                    row.append("\n".join(texts).strip())
                row = _trim_trailing_empty(row)
                if _row_has_text(row):
                    rows.append(row)
            if rows:
                out.append(
                    ExtractedTable(
                        source=source,
                        name=f"{Path(name).name} table {len(out) + 1}",
                        rows=rows,
                        table_index=len(out) + 1,
                    )
                )
    return out


def _hwpx_bytes_has_table(hwpx_bytes: bytes) -> bool:
    """HWPX(zip) 어느 XML에든 표 구조(<...:tbl>)가 있으면 True (존재만 판정)."""
    import io
    import re as _re
    import zipfile

    try:
        z = zipfile.ZipFile(io.BytesIO(hwpx_bytes))
        _guard_zip_bomb(z, source="hwpx")
    except Exception:  # noqa: BLE001
        return False
    for n in z.namelist():
        if not n.lower().endswith(".xml"):
            continue
        try:
            xml = z.read(n).decode("utf-8", errors="replace")
        except Exception:  # noqa: BLE001
            continue
        if _re.search(r"<(?:\w+:)?tbl\b", xml):
            return True
    return False


def _markdown_tables(md: str, *, source: str, name_prefix: str) -> list[ExtractedTable]:
    """마크다운 표(``| a | b |``)를 ExtractedTable 로 파싱.

    연속된 파이프 행을 하나의 표로 묶고, 구분선(``| --- |``)과 전 셀이 빈 행은 버린다.
    셀 안 줄바꿈은 unhwp 가 ``<br>`` 로 내보내므로 공백으로 환원한다.
    """
    import re as _re

    tables: list[ExtractedTable] = []
    block: list[list[str]] = []

    def flush() -> None:
        if block:
            tables.append(
                ExtractedTable(
                    source=source,
                    name=f"{name_prefix} table {len(tables) + 1}",
                    rows=list(block),
                    table_index=len(tables) + 1,
                )
            )
            block.clear()

    for line in (md or "").splitlines():
        s = line.strip()
        if not s.startswith("|"):
            flush()
            continue
        # 구분선(--- / :---:)은 구조일 뿐 내용이 아니다.
        if not _re.sub(r"[|\s:-]", "", s):
            continue
        cells = [_cell_text(_re.sub(r"<br\s*/?>", " ", c)) for c in s.strip("|").split("|")]
        cells = _trim_trailing_empty(cells)
        if _row_has_text(cells):
            block.append(cells)
    flush()
    return tables


def _hwp_tables_via_unhwp(p: Path) -> tuple[str | None, list[ExtractedTable]]:
    """unhwp(MIT)로 .hwp 표 셀을 회수 — 표 부분만 돌려준다(본문 대체 아님).

    rhwp 는 .hwp 바이너리에서 표 셀을 통째로 흘린다(실측: 공고문 표 0행). unhwp 의
    HWP5 파서는 표를 마크다운으로 렌더하므로 그 표 행만 뽑아, 호출부가 본문에
    **합집합**으로 덧붙인다.

    **.hwp 전용이다.** .hwpx 는 rhwp + _hwpx_tables 가 이미 정확한 반면 unhwp 는
    표를 통째로 누락한 실측 사례가 있다(acc-S1-03.hwpx: 표 5행→0행, 인수 기대숫자
    3/3→0/3). 형식을 넓히면 회귀한다.

    라이선스: unhwp 는 MIT. 네이티브 부분은 Rust 이고 동봉 크레이트 29종 전수가
    MIT/Apache-2.0/Unlicense/Zlib — 카피레프트도, 기존 AGPL 구현(pyhwp/hwp.js/
    hwplib) 파생도 없음을 바이너리 문자열 + crates.io 조회로 확인(2026-08-01).
    구 pyhwp(AGPL, hwp5html CLI arm's length) 경로를 대체한다.

    미설치·파싱 실패는 (None, []) — 호출부가 rhwp 결과를 그대로 쓴다.
    """
    if p.suffix.lower() != ".hwp":
        return None, []
    try:
        import unhwp  # type: ignore
    except ImportError:
        return None, []
    try:
        md = unhwp.to_markdown(str(p))
    except Exception:  # noqa: BLE001
        return None, []
    if not md:
        return None, []
    tables = _markdown_tables(md, source="hwp", name_prefix=p.name)
    if not tables:
        return None, []
    return _tables_to_text(tables), tables


def _docx_collect_tables(container, *, source: str, tables: list[ExtractedTable],
                         parts: list[str]) -> None:
    for idx, tbl in enumerate(getattr(container, "tables", []), start=1):
        rows: list[list[str]] = []
        for row in tbl.rows:
            cells = _trim_trailing_empty([_cell_text(cell.text) for cell in row.cells])
            if _row_has_text(cells):
                rows.append(cells)
                parts.append(_table_text_line(cells))
        if rows:
            tables.append(
                ExtractedTable(
                    source=source,
                    name=f"{source} table {idx}",
                    rows=rows,
                    table_index=idx,
                )
            )


def _docx_ooxml_extras(p: Path, warnings: list[str]) -> list[str]:
    import re as _re
    import zipfile

    extras: list[str] = []
    names = _zip_names(p)
    if any(n.startswith("word/media/") for n in names):
        _warn_once(warnings, "docx_media_not_ocrd")
        extras.extend(_ocr_package_images(p, prefix="word/media/", label="docx", warnings=warnings))
    if any(n.startswith("word/charts/") for n in names):
        _warn_once(warnings, "docx_charts_not_extracted")
    if any(n.startswith("word/embeddings/") for n in names):
        _warn_once(warnings, "docx_embedded_objects_not_extracted")

    try:
        with zipfile.ZipFile(p) as zf:
            _guard_zip_bomb(zf, source="office-zip")
            named_parts = {
                "word/comments.xml": "docx_comments_extracted",
                "word/footnotes.xml": "docx_footnotes_extracted",
                "word/endnotes.xml": "docx_endnotes_extracted",
            }
            for name, warning in named_parts.items():
                if name not in names:
                    continue
                text = _xml_text(zf.read(name).decode("utf-8", errors="replace"))
                if text:
                    extras.append(f"[{name}]\n{text}")
                    _warn_once(warnings, warning)

            xml_parts = [
                n for n in names
                if n.startswith("word/") and n.endswith(".xml")
                and (n == "word/document.xml" or n.startswith("word/header") or n.startswith("word/footer"))
            ]
            for name in xml_parts:
                xml = zf.read(name).decode("utf-8", errors="replace")
                for box in _re.findall(
                    r"<(?:\w+:)?txbxContent\b[^>]*>(.*?)</(?:\w+:)?txbxContent>",
                    xml,
                    _re.S,
                ):
                    text = _xml_text(box)
                    if text:
                        extras.append(f"[{name}:textbox]\n{text}")
                        _warn_once(warnings, "docx_textboxes_extracted")
    except Exception:  # noqa: BLE001
        return extras
    return extras


def _docx_ooxml_tables(p: Path, warnings: list[str]) -> list[ExtractedTable]:
    import re as _re
    import zipfile

    tables: list[ExtractedTable] = []
    names = _zip_names(p)
    xml_parts = [
        n for n in names
        if n.startswith("word/") and n.endswith(".xml")
        and (n == "word/document.xml" or n.startswith("word/header") or n.startswith("word/footer"))
    ]
    try:
        with zipfile.ZipFile(p) as zf:
            _guard_zip_bomb(zf, source="office-zip")
            for name in xml_parts:
                xml = zf.read(name).decode("utf-8", errors="replace")
                for raw_tbl in _re.findall(
                    r"<(?:\w+:)?tbl\b[^>]*>(.*?)</(?:\w+:)?tbl>",
                    xml,
                    _re.S,
                ):
                    rows: list[list[str]] = []
                    for raw_tr in _re.findall(
                        r"<(?:\w+:)?tr\b[^>]*>(.*?)</(?:\w+:)?tr>",
                        raw_tbl,
                        _re.S,
                    ):
                        row: list[str] = []
                        for raw_tc in _re.findall(
                            r"<(?:\w+:)?tc\b[^>]*>(.*?)</(?:\w+:)?tc>",
                            raw_tr,
                            _re.S,
                        ):
                            row.append(_xml_text(raw_tc))
                        row = _trim_trailing_empty(row)
                        if _row_has_text(row):
                            rows.append(row)
                    if rows:
                        tables.append(
                            ExtractedTable(
                                source="docx_ooxml",
                                name=f"{Path(name).name} table {len(tables) + 1}",
                                rows=rows,
                                table_index=len(tables) + 1,
                            )
                        )
    except Exception:  # noqa: BLE001
        return tables
    if tables:
        _warn_once(warnings, "docx_ooxml_tables_extracted")
    return tables


def _docx_ooxml_text(p: Path, warnings: list[str]) -> str:
    import zipfile

    names = _zip_names(p)
    xml_parts = [
        n for n in names
        if n.startswith("word/") and n.endswith(".xml")
        and (n == "word/document.xml" or n.startswith("word/header") or n.startswith("word/footer"))
    ]
    lines: list[str] = []
    try:
        with zipfile.ZipFile(p) as zf:
            _guard_zip_bomb(zf, source="office-zip")
            for name in xml_parts:
                text = _xml_text(zf.read(name).decode("utf-8", errors="replace"))
                if text:
                    lines.extend(text.splitlines())
    except Exception:  # noqa: BLE001
        return ""
    text = "\n".join(line for line in lines if line.strip())
    return text


def _extract_docx(p: Path) -> ExtractResult:
    try:
        from docx import Document  # type: ignore

        doc = Document(str(p))
        parts = [para.text for para in doc.paragraphs if para.text]
        tables: list[ExtractedTable] = []
        warnings: list[str] = []

        _docx_collect_tables(doc, source="docx", tables=tables, parts=parts)
        for sidx, section in enumerate(doc.sections, start=1):
            for label in (
                "header",
                "footer",
                "first_page_header",
                "first_page_footer",
                "even_page_header",
                "even_page_footer",
            ):
                block = getattr(section, label, None)
                if block is None:
                    continue
                block_parts = [p.text for p in block.paragraphs if p.text]
                if block_parts:
                    parts.append(f"[docx section {sidx} {label}]")
                    parts.extend(block_parts)
                _docx_collect_tables(
                    block,
                    source=f"docx_{label}",
                    tables=tables,
                    parts=parts,
                )

        if not tables:
            ooxml_tables = _docx_ooxml_tables(p, warnings)
            if ooxml_tables:
                tables.extend(ooxml_tables)
                table_text = _tables_to_text(ooxml_tables)
                merged = _append_if_missing("\n".join(parts), table_text)
                parts = merged.splitlines() if merged else []

        parts.extend(_docx_ooxml_extras(p, warnings))
        ooxml_text = _docx_ooxml_text(p, warnings)
        if ooxml_text:
            current = "\n".join(parts)
            merged = _append_if_missing(current, ooxml_text)
            if merged != current:
                _warn_once(warnings, "docx_ooxml_text_extracted")
            parts = merged.splitlines() if merged else []
        return ExtractResult(
            text="\n".join(parts),
            method="parser",
            quality=0.97,
            tables=tables,
            warnings=warnings,
        )
    except Exception as exc:  # noqa: BLE001
        return ExtractResult(text="", method="parser", quality=0.0, error=str(exc))


def _excel_package_warnings(p: Path, warnings: list[str]) -> None:
    names = _zip_names(p)
    if any(n.startswith("xl/media/") for n in names):
        _warn_once(warnings, "excel_media_not_ocrd")
    if any(n.startswith("xl/charts/") for n in names):
        _warn_once(warnings, "excel_charts_not_extracted")
    if any(n.startswith("xl/drawings/") for n in names):
        _warn_once(warnings, "excel_drawings_text_may_be_missing")
    if any(n.startswith("xl/embeddings/") for n in names):
        _warn_once(warnings, "excel_embedded_objects_not_extracted")


def _excel_auxiliary_texts(path: str, warnings: list[str]) -> list[str]:
    parts: list[str] = []
    try:
        from openpyxl import load_workbook  # type: ignore

        wb = load_workbook(path, read_only=False, data_only=False)
    except Exception:  # noqa: BLE001
        return parts
    try:
        for ws in wb.worksheets:
            for row in ws.iter_rows():
                for cell in row:
                    coord = f"{ws.title}!{cell.coordinate}"
                    if cell.comment and cell.comment.text:
                        parts.append(f"[excel comment {coord}] {cell.comment.text.strip()}")
                        _warn_once(warnings, "excel_comments_extracted")
                    if cell.hyperlink and cell.hyperlink.target:
                        parts.append(f"[excel hyperlink {coord}] {cell.hyperlink.target}")
                        _warn_once(warnings, "excel_hyperlinks_extracted")
                    if isinstance(cell.value, str) and cell.value.startswith("="):
                        parts.append(f"[excel formula {coord}] {cell.value}")
                        _warn_once(warnings, "excel_formulas_extracted")
    except Exception:  # noqa: BLE001
        return parts
    finally:
        wb.close()
    return parts


def _extract_excel(p: Path) -> ExtractResult:
    """Excel(.xlsx/.xlsm/.xls) → 시트별 행을 텍스트로.

    표 구조는 셀을 ' | '로 잇고 시트마다 '[sheet: 이름]' 헤더를 붙인다(DOCX 표와 동일 방식).
    .xlsx/.xlsm은 openpyxl(read_only·data_only=수식 대신 값), .xls는 xlrd(설치 시).
    라이브러리 미설치/파싱 실패는 graceful degrade — 빈 텍스트 + error.

    데이터 누락 가드(2026-06):
      - 숨김(hidden/veryHidden) 시트도 포함한다. 숨김 시트에 비밀 수치가 들어 있을 수
        있으므로 sheet_state로 거르지 않고 모두 추출하되 헤더에 상태를 표기한다.
      - data_only=True는 '캐시된 수식 값'을 읽는다. 엑셀이 한 번도 계산·저장하지 않은
        수식 셀은 캐시가 없어 None으로 나와 누락될 수 있다. None인 수식 셀이 감지되면
        error에 경고를 남겨 후속 검수가 LibreOffice 재계산 등을 판단하게 한다.
      - 병합 셀(merged): read_only는 병합정보를 노출하지 않아, 세로로 병합된 분류/카테고리
        라벨이 첫 행에만 남고 나머지 행에서 끊긴다. _excel_merged_values로 별도 패스를 떠서
        각 병합 범위의 앵커 값을 '걸친 모든 행의 맨 왼쪽 열'에 복제한다(가로 병합은 앵커 한
        칸만 유지해 중복 노이즈 회피).
    """
    suffix = p.suffix.lower().lstrip(".")
    warnings: list[str] = []
    if suffix in ("xlsx", "xlsm"):
        _excel_package_warnings(p, warnings)
        try:
            from openpyxl import load_workbook  # type: ignore
        except ImportError as exc:
            return ExtractResult(text="", method="openpyxl", quality=0.0,
                                 error=f"openpyxl not installed: {exc}")
        try:
            wb = load_workbook(str(p), read_only=True, data_only=True)
        except Exception as exc:  # noqa: BLE001
            return ExtractResult(text="", method="openpyxl", quality=0.0, error=str(exc))
        merged = _excel_merged_values(str(p))
        parts: list[str] = []
        tables: list[ExtractedTable] = []
        # read_only 워크북도 hidden/veryHidden 시트를 worksheets에 노출하므로
        # 추가 필터 없이 전부 순회 — sheet_state는 헤더 표기로만 사용.
        for ws_idx, ws in enumerate(wb.worksheets, start=1):
            state = getattr(ws, "sheet_state", "visible")
            header = f"[sheet: {ws.title}]" if state == "visible" else f"[sheet: {ws.title} ({state})]"
            parts.append(header)
            mmap = merged.get(ws.title)
            table_rows: list[list[str]] = []
            if not mmap:
                # 병합 없는 시트는 기존 fast path(values_only) 유지.
                for row in ws.iter_rows(values_only=True):
                    cells = _trim_trailing_empty([_cell_text(c) for c in row])
                    if _row_has_text(cells):
                        table_rows.append(cells)
                        parts.append(_table_text_line(cells))
                if table_rows:
                    tables.append(
                        ExtractedTable(
                            source="excel",
                            name=ws.title,
                            rows=table_rows,
                            sheet=ws.title,
                            table_index=ws_idx,
                        )
                    )
                continue
            # 병합 있는 시트만 좌표를 추적해 None 셀에 전파값을 채운다.
            # read_only는 빈 칸을 EmptyCell(좌표 속성 없음)로 주므로, 행마다 첫
            # 실셀의 (row, column)에서 행번호·열 오프셋을 잡아 위치를 보정한다.
            for row in ws.iter_rows():
                row_no = col_base = None
                for idx, cell in enumerate(row):
                    if hasattr(cell, "column"):
                        row_no, col_base = cell.row, cell.column - idx
                        break
                cells = []
                for idx, cell in enumerate(row):
                    v = cell.value
                    if v is None and row_no is not None:
                        v = mmap.get((row_no, col_base + idx))
                    cells.append(_cell_text(v))
                cells = _trim_trailing_empty(cells)
                if _row_has_text(cells):
                    table_rows.append(cells)
                    parts.append(_table_text_line(cells))
            if table_rows:
                tables.append(
                    ExtractedTable(
                        source="excel",
                        name=ws.title,
                        rows=table_rows,
                        sheet=ws.title,
                        table_index=ws_idx,
                    )
                )
        wb.close()

        # 수식 캐시 부재(None) 가능성 감지 — data_only 모드에서 수식 셀이 None이면
        # 값 누락. read_only 워크북은 셀 데이터타입을 못 보므로 별도 워크북을
        # data_only=False로 재오픈해 수식 존재 여부만 가볍게 점검한다(실패는 무시).
        parts.extend(_excel_auxiliary_texts(str(p), warnings))
        parts.extend(_ocr_package_images(p, prefix="xl/media/", label="excel", warnings=warnings))
        warn = _excel_formula_cache_warning(str(p))
        if warn:
            _warn_once(warnings, "excel_formula_cached_values_missing")
        return ExtractResult(
            text="\n".join(parts),
            method="openpyxl",
            quality=0.95,
            error=warn,
            tables=tables,
            warnings=warnings,
        )

    # .xls (구형 바이너리) — openpyxl 미지원이라 xlrd 필요. 미설치 시 변환 안내.
    try:
        import xlrd  # type: ignore
    except ImportError as exc:
        return ExtractResult(text="", method="xlrd", quality=0.0,
                             error=f".xls requires xlrd (or convert to .xlsx): {exc}")
    try:
        book = xlrd.open_workbook(str(p))
    except Exception as exc:  # noqa: BLE001
        return ExtractResult(text="", method="xlrd", quality=0.0, error=str(exc))
    parts = []
    tables: list[ExtractedTable] = []
    for sidx, sh in enumerate(book.sheets(), start=1):
        parts.append(f"[sheet: {sh.name}]")
        table_rows: list[list[str]] = []
        for r in range(sh.nrows):
            cells = _trim_trailing_empty([
                _xls_cell_text(sh.cell_value(r, c), sh.cell_type(r, c), book.datemode)
                for c in range(sh.ncols)
            ])
            if _row_has_text(cells):
                table_rows.append(cells)
                parts.append(_table_text_line(cells))
        if table_rows:
            tables.append(
                ExtractedTable(
                    source="excel",
                    name=sh.name,
                    rows=table_rows,
                    sheet=sh.name,
                    table_index=sidx,
                )
            )
    return ExtractResult(text="\n".join(parts), method="xlrd", quality=0.9, tables=tables)


# 병합 셀 세로 전파 상한 — 병리적으로 큰 병합(예: A1:A1048576)이 맵을 폭발시키지
# 않도록 범위당 전파 행 수에 상한. 초과 시 그 범위는 앵커(좌상단)만 자연 유지.
_MERGED_MAX_ROWS = 10000


def _excel_merged_values(path: str) -> dict[str, dict[tuple[int, int], object]]:
    """시트별 병합 셀의 '세로 라벨 전파' 맵 (best-effort).

    read_only 모드는 병합정보를 노출하지 않으므로 비-read_only로 한 번 더 열어
    각 병합 범위의 앵커(좌상단) 값을 범위가 걸친 모든 행의 '맨 왼쪽 열'에 복제한다.
    세로로 병합된 분류/카테고리 라벨이 각 행에 함께 남아 표 의미가 끊기지 않는다.
    가로 병합은 앵커 한 칸만 남겨 중복 노이즈를 피한다(원동작 유지).

    반환: {sheet_title: {(row, min_col): anchor_value}}. 라이브러리/파싱 실패는
    조용히 빈 맵 — 전파 없이도 본문 추출은 정상 동작.
    """
    out: dict[str, dict[tuple[int, int], object]] = {}
    try:
        from openpyxl import load_workbook  # type: ignore

        wb = load_workbook(path, data_only=True)  # 병합정보 필요 → read_only 불가
    except Exception:  # noqa: BLE001
        return out
    try:
        for ws in wb.worksheets:
            ranges = getattr(ws, "merged_cells", None)
            if not ranges:
                continue
            m: dict[tuple[int, int], object] = {}
            for rng in list(getattr(ranges, "ranges", [])):
                if rng.max_row - rng.min_row + 1 > _MERGED_MAX_ROWS:
                    continue
                anchor = ws.cell(row=rng.min_row, column=rng.min_col).value
                if anchor is None or not str(anchor).strip():
                    continue
                for r in range(rng.min_row, rng.max_row + 1):
                    m[(r, rng.min_col)] = anchor
            if m:
                out[ws.title] = m
    except Exception:  # noqa: BLE001
        return out
    finally:
        wb.close()
    return out


def _excel_formula_cache_warning(path: str) -> str | None:
    """data_only=True에서 수식 캐시 부재로 값이 None일 위험을 가볍게 점검.

    data_only=False로 재오픈해 수식 셀(value가 '='로 시작) 중 별도 캐시 워크북에서
    None인 게 있는지 본다. 라이브러리/파싱 실패는 조용히 None(경고 없음) — best-effort.
    LibreOffice 재계산 같은 본격 복구는 범위 밖(risks 참조).
    """
    try:
        from openpyxl import load_workbook  # type: ignore

        wb_f = load_workbook(path, read_only=True, data_only=False)
        wb_v = load_workbook(path, read_only=True, data_only=True)
    except Exception:  # noqa: BLE001
        return None
    missing = 0
    try:
        for ws_f, ws_v in zip(wb_f.worksheets, wb_v.worksheets):
            for row_f, row_v in zip(ws_f.iter_rows(), ws_v.iter_rows()):
                for cf, cv in zip(row_f, row_v):
                    fv = cf.value
                    # 스칼라 수식은 str('='시작), 배열/데이터테이블 수식은 openpyxl이 ArrayFormula·
                    # DataTableFormula 객체로 준다(str 아님) — isinstance(str) 만 보면 배열수식의
                    # 캐시 부재를 놓쳐 계산값이 조용히 유실된다. 타입명으로 판별(구버전 import 안전).
                    is_scalar_formula = isinstance(fv, str) and fv.startswith("=")
                    is_formula_obj = type(fv).__name__ in ("ArrayFormula", "DataTableFormula")
                    if (is_scalar_formula or is_formula_obj) and cv.value is None:
                        missing += 1
    except Exception:  # noqa: BLE001
        return None
    finally:
        wb_f.close()
        wb_v.close()
    if missing:
        return (
            f"warning: {missing} formula cell(s) had no cached value (data_only) — "
            "값이 누락됐을 수 있음. LibreOffice 등으로 재계산 후 재추출 권장."
        )
    return None


def _extract_doc(p: Path) -> ExtractResult:
    """구형 .doc(Word 97-2003) → antiword CLI 추출.

    python-docx는 .docx(OOXML) 전용이라 .doc 바이너리는 못 읽는다. 시스템에 antiword가
    있으면 그것으로 텍스트 추출, 없으면 graceful degrade(.docx 변환 안내).
    한국어는 antiword 매핑 의존 — UTF-8 매핑(-m UTF-8.txt) 우선 시도 후 기본.
    """
    import shutil
    import subprocess

    exe = shutil.which("antiword")
    if not exe:
        return ExtractResult(
            text="", method="antiword", quality=0.0,
            error=".doc 추출에는 antiword 필요 (또는 .docx로 변환)",
        )
    out = None
    for args in ([exe, "-m", "UTF-8.txt", str(p)], [exe, str(p)]):
        try:
            out = subprocess.run(args, capture_output=True, timeout=60)
        except Exception as exc:  # noqa: BLE001
            return ExtractResult(text="", method="antiword", quality=0.0, error=str(exc))
        if out.returncode == 0 and out.stdout.strip():
            return ExtractResult(
                text=out.stdout.decode("utf-8", errors="replace"),
                method="antiword", quality=0.9,
            )
    err = (out.stderr.decode("utf-8", errors="replace") if out and out.stderr else "")[:200]
    return ExtractResult(text="", method="antiword", quality=0.0,
                         error=err or "antiword: no text extracted")


def _extract_pptx(p: Path) -> ExtractResult:
    """PowerPoint(.pptx/.pptm) → 슬라이드별 도형 텍스트 + 표.

    슬라이드마다 '[slide N]' 헤더. 텍스트 도형·표 셀을 추출(이미지/차트 제외).
    .ppt(구형 바이너리)는 python-pptx 미지원 — 미지원 처리(변환 권장).
    """
    try:
        from pptx import Presentation  # type: ignore
    except ImportError as exc:
        return ExtractResult(text="", method="pptx", quality=0.0,
                             error=f"python-pptx not installed: {exc}")
    try:
        prs = Presentation(str(p))
    except Exception as exc:  # noqa: BLE001
        return ExtractResult(text="", method="pptx", quality=0.0, error=str(exc))
    parts: list[str] = []
    tables: list[ExtractedTable] = []
    warnings: list[str] = []
    for i, slide in enumerate(prs.slides):
        parts.append(f"[slide {i + 1}]")
        for shape_idx, shape in enumerate(slide.shapes, start=1):
            if getattr(shape, "has_text_frame", False):
                for para in shape.text_frame.paragraphs:
                    t = "".join(r.text for r in para.runs) or para.text
                    if t.strip():
                        parts.append(t)
            if getattr(shape, "has_table", False):
                table_rows: list[list[str]] = []
                for row in shape.table.rows:
                    cells = _trim_trailing_empty([_cell_text(c.text) for c in row.cells])
                    if _row_has_text(cells):
                        table_rows.append(cells)
                        parts.append(_table_text_line(cells))
                if table_rows:
                    tables.append(
                        ExtractedTable(
                            source="pptx",
                            name=f"slide {i + 1} table {shape_idx}",
                            rows=table_rows,
                            page=i + 1,
                            table_index=shape_idx,
                        )
                    )
            if getattr(shape, "has_chart", False):
                _warn_once(warnings, "pptx_charts_not_extracted")
            shape_type = str(getattr(shape, "shape_type", "")).upper()
            if "PICTURE" in shape_type:
                _warn_once(warnings, "pptx_media_not_ocrd")
            if "OLE" in shape_type or "EMBEDDED" in shape_type:
                _warn_once(warnings, "pptx_embedded_objects_not_extracted")
        if getattr(slide, "has_notes_slide", False):
            try:
                note_text = slide.notes_slide.notes_text_frame.text
                if note_text and note_text.strip():
                    parts.append(f"[slide {i + 1} notes]")
                    parts.append(note_text.strip())
                    _warn_once(warnings, "pptx_notes_extracted")
            except Exception:  # noqa: BLE001
                pass
    parts.extend(_ocr_package_images(p, prefix="ppt/media/", label="pptx", warnings=warnings))
    return ExtractResult(
        text="\n".join(parts),
        method="pptx",
        quality=0.95,
        tables=tables,
        warnings=warnings,
    )


def _pdf_tables_via_pdfplumber(p: Path) -> tuple[list[ExtractedTable], list[str]]:
    try:
        import pdfplumber  # type: ignore
    except ImportError:
        return [], ["pdf_table_extractor_unavailable"]
    except Exception as exc:  # noqa: BLE001
        return [], [f"pdf_table_extractor_error:{type(exc).__name__}"]

    tables: list[ExtractedTable] = []
    try:
        with pdfplumber.open(str(p)) as pdf:
            for page_idx, page in enumerate(pdf.pages, start=1):
                for raw in page.extract_tables() or []:
                    rows: list[list[str]] = []
                    for row in raw or []:
                        cells = _trim_trailing_empty([_cell_text(c) for c in (row or [])])
                        if _row_has_text(cells):
                            rows.append(cells)
                    if rows:
                        tables.append(
                            ExtractedTable(
                                source="pdf",
                                name=f"page {page_idx} table {len(tables) + 1}",
                                rows=rows,
                                page=page_idx,
                                table_index=len(tables) + 1,
                            )
                        )
    except Exception as exc:  # noqa: BLE001
        return [], [f"pdf_table_extractor_error:{type(exc).__name__}"]

    if tables:
        return tables, ["pdf_tables_structured"]
    return [], ["pdf_tables_not_detected"]


def _pdf_table_coverage(table_warnings: list[str]) -> str | None:
    if any(w.startswith("pdf_table_extractor_") for w in table_warnings):
        return "incomplete"
    return None


def _pdf_sparse_signal(text: str, pages: int, base_quality: float) -> tuple[float, list[str]]:
    """부분 텍스트레이어(표지/워터마크/머리말)가 스캔 본문의 OCR 을 가로채는 것 방지(#10).

    pdfminer/PyMuPDF 는 텍스트가 조금이라도 있으면 즉시 채택하고 OCR 로 안 내려간다. 다수 페이지에
    문자밀도가 매우 낮으면(준스캔 의심) quality 를 검수 임계(extraction_review_min_quality 기본 0.6)
    아래로 낮춰 자동확정 대신 검수 라우팅되게 한다(table_coverage 와 동형 FNR-safe 신호, 등급·본문
    불변). 정상 밀도면 base_quality 유지. 순수 스캔(레이어 0)은 이 분기 전에 OCR 로 처리된다.
    """
    density = len(text.strip()) / max(1, pages)
    if pages >= 3 and density < 80:
        return 0.5, [f"pdf_sparse_text_layer: {density:.0f} chars/page over {pages}p (준스캔 의심 → 검수)"]
    return base_quality, []


def _extract_pdf(p: Path) -> ExtractResult:
    """PDF 추출.

    1) pdfminer.six (BSD-3) — 텍스트 레이어 있는 경우
    2) PyMuPDF (AGPL) — 설치된 경우 fallback
    3) Tesseract OCR — 텍스트 레이어 없는 스캔본
    """
    # 1순위: pdfminer.six
    try:
        from pdfminer.high_level import extract_text  # type: ignore

        text = extract_text(str(p)) or ""
        if text.strip():
            pages = max(1, text.count("\x0c") + 1)
            tables, table_warnings = _pdf_tables_via_pdfplumber(p)
            text = _append_if_missing(text, _tables_to_text(tables))
            q, sparse_w = _pdf_sparse_signal(text, pages, 0.92)
            return ExtractResult(
                text=text,
                method="parser",
                quality=q,
                pages=pages,
                table_coverage=_pdf_table_coverage(table_warnings),
                tables=tables,
                warnings=table_warnings + sparse_w,
            )
    except Exception:  # noqa: BLE001
        pass

    # 2순위: PyMuPDF (AGPL 옵트인)
    try:
        import fitz  # type: ignore

        doc = fitz.open(str(p))
        try:
            page_texts = [page.get_text() for page in doc]
            n_pages = len(page_texts)
        finally:
            doc.close()
        text = "\n".join(page_texts)
        if text.strip():
            tables, table_warnings = _pdf_tables_via_pdfplumber(p)
            text = _append_if_missing(text, _tables_to_text(tables))
            q, sparse_w = _pdf_sparse_signal(text, n_pages, 0.95)
            return ExtractResult(
                text=text,
                method="parser",
                quality=q,
                pages=n_pages,
                table_coverage=_pdf_table_coverage(table_warnings),
                tables=tables,
                warnings=table_warnings + sparse_w,
            )
        # 텍스트 레이어 없음 → OCR 시도
        return _ocr_pdf_pages(p, n_pages)
    except Exception:  # noqa: BLE001
        pass

    # 3순위: pdfminer로 읽혔지만 빈 텍스트 → OCR 시도
    return _ocr_pdf_pages(p, None)


def _ocr_pdf_pages(p: Path, n_pages: int | None) -> ExtractResult:
    """PDF 페이지를 이미지로 변환 후 Tesseract OCR.

    DoS 가드: convert_from_path를 first_page/last_page로 [1, max_pages] 범위만
    렌더한다. 상한을 넘는 페이지는 변환·OCR하지 않고 경고를 error에 남긴다(잘림 표기).
    상한이 0 이하면 무제한(opt-out).
    """
    try:
        from pdf2image import convert_from_path  # type: ignore

        kwargs: dict = {"dpi": 200}
        if POPPLER_PATH:
            kwargs["poppler_path"] = POPPLER_PATH

        max_pages = _ocr_max_pages()
        truncated = False
        last_page: int | None = None
        if max_pages and max_pages > 0:
            # n_pages를 아는 경우 초과 여부를 미리 판정, 모르면 상한까지만 렌더.
            last_page = max_pages
            if n_pages is not None and n_pages > max_pages:
                truncated = True
            kwargs["first_page"] = 1
            kwargs["last_page"] = last_page

        images = convert_from_path(str(p), **kwargs)
        # n_pages를 몰랐어도 렌더 결과가 상한과 같으면 잘렸을 수 있음(보수적 표기).
        if last_page is not None and len(images) >= last_page and n_pages is None:
            truncated = True
        texts = [_tess_image(img) for img in images]
        text = "\n".join(t for t in texts if t)
        if text.strip():
            note = (
                f"OCR truncated to first {last_page} of {n_pages or '?'} pages (DoS guard)"
                if truncated else None
            )
            warnings = ["pdf_ocr_text_may_be_noisy"]
            if truncated:
                warnings.append("pdf_ocr_truncated")
            return ExtractResult(
                text=text, method="ocr", quality=0.75,
                ocr_used=True, pages=len(images), error=note, warnings=warnings,
            )
    except Exception:  # noqa: BLE001
        pass
    return ExtractResult(
        text="", method="ocr", quality=0.0, ocr_used=False,
        pages=n_pages,
        error="OCR failed (no text layer, pdf2image/tesseract unavailable)",
        warnings=["pdf_ocr_unavailable"],
    )


def _extract_image_ocr(p: Path) -> ExtractResult:
    """이미지 파일(jpg/png/tiff 등) 직접 OCR."""
    try:
        from PIL import Image  # type: ignore

        img = Image.open(str(p))
        text = _tess_image(img)
        if text.strip():
            return ExtractResult(
                text=text,
                method="ocr",
                quality=0.75,
                ocr_used=True,
                warnings=["image_ocr_text_may_be_noisy"],
            )
        return ExtractResult(text="", method="ocr", quality=0.0, ocr_used=True,
                             error="OCR produced empty text")
    except Exception as exc:  # noqa: BLE001
        return ExtractResult(text="", method="ocr", quality=0.0,
                             ocr_used=True, error=str(exc))


def _tess_image(img) -> str:
    """PIL Image → Tesseract OCR 텍스트. kor+eng 병행 인식."""
    import pytesseract  # type: ignore

    pytesseract.pytesseract.tesseract_cmd = TESSERACT_CMD
    langs = "+".join(_available_tess_langs(["kor", "eng"]))
    return pytesseract.image_to_string(img, lang=langs, config="--psm 3")


def _available_tess_langs(preferred: list[str]) -> list[str]:
    """설치된 언어팩 중 preferred 교집합. 없으면 eng 폴백."""
    try:
        import pytesseract  # type: ignore

        pytesseract.pytesseract.tesseract_cmd = TESSERACT_CMD
        installed = set(pytesseract.get_languages())
        result = [lg for lg in preferred if lg in installed]
        return result if result else ["eng"]
    except Exception:  # noqa: BLE001
        return ["eng"]
